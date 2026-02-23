#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
创建平台经济高质量发展PPT
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import datetime
import os

def create_platform_economy_ppt():
    """创建平台经济高质量发展PPT"""
    
    # 创建演示文稿
    prs = Presentation()
    
    # 设置幻灯片宽度和高度（16:9）
    prs.slide_width = Inches(13.33)  # 25.4 cm
    prs.slide_height = Inches(7.5)   # 19.05 cm
    
    # ========== 1. 封面页 ==========
    slide_layout = prs.slide_layouts[0]  # 标题幻灯片
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "平台经济高质量发展\n路径与对策"
    subtitle.text = f"数字经济时代下的新引擎\n{datetime.datetime.now().strftime('%Y年%m月%d日')}"
    
    # 设置标题样式
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)  # 深蓝色
    
    # ========== 2. 目录页 ==========
    slide_layout = prs.slide_layouts[1]  # 标题和内容
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "目录"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    content = slide.shapes.placeholders[1]
    tf = content.text_frame
    
    # 清空默认文本
    tf.clear()
    
    # 添加目录项
    items = [
        "一、平台经济的内涵与特征",
        "二、平台经济发展现状与趋势", 
        "三、平台经济发展面临的挑战",
        "四、平台经济高质量发展对策",
        "五、未来发展方向与展望"
    ]
    
    for item in items:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(24)
        p.font.bold = True
        p.space_after = Pt(12)
        p.level = 0
        
    # ========== 3. 内涵与特征 ==========
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "一、平台经济的内涵与特征"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    content = slide.shapes.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    sections = [
        ("📊 内涵定义", [
            "• 平台经济是以互联网平台为主要载体",
            "• 以数据为关键生产要素",
            "• 以新一代信息技术为核心驱动力",
            "• 以网络信息基础设施为重要支撑的新型经济形态"
        ]),
        ("🎯 基本特征", [
            "• 虚拟性：摆脱时间与空间约束",
            "• 网络外部性：用户越多，价值越大",
            "• 多元归属性：用户可以参与多个平台",
            "• 创新协同性：数实融合的重要桥梁"
        ]),
        ("🏢 四个层面", [
            "• 数字平台：引擎",
            "• 平台企业：主体", 
            "• 平台生态：载体",
            "• 平台经济：整体"
        ])
    ]
    
    for section_title, points in sections:
        p = tf.add_paragraph()
        p.text = section_title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 102, 204)
        p.space_after = Pt(6)
        
        for point in points:
            p = tf.add_paragraph()
            p.text = point
            p.font.size = Pt(18)
            p.level = 1
            p.space_before = Pt(3)
    
    # ========== 4. 发展现状 ==========
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "二、平台经济发展现状与趋势"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    content = slide.shapes.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    sections = [
        ("📈 发展历程", [
            "• 2015年：国外文献出现平台经济概念",
            "• 2018年：首次写入《政府工作报告》",
            "• 2021年：《关于推动平台经济规范健康持续发展的若干意见》发布",
            "• 2023年：中央经济工作会议强调促进平台经济发展"
        ]),
        ("🎪 主要业态", [
            "• 网络销售：电商平台",
            "• 生活服务：外卖、出行",
            "• 社交娱乐：社交媒体、短视频",
            "• 信息服务：资讯、搜索",
            "• 金融服务：移动支付、金融科技",
            "• 产业平台：工业互联网"
        ]),
        ("🚀 发展趋势", [
            "• 从消费互联网向产业互联网过渡",
            "• 平台规模扩张走向生态化",
            "• 中小企业数字化转型成为重要场景",
            "• 与产业链供应链深度融合"
        ])
    ]
    
    for section_title, points in sections:
        p = tf.add_paragraph()
        p.text = section_title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 102, 204)
        p.space_after = Pt(6)
        
        for point in points:
            p = tf.add_paragraph()
            p.text = point
            p.font.size = Pt(18)
            p.level = 1
            p.space_before = Pt(3)
    
    # ========== 5. 面临的挑战 ==========
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "三、平台经济发展面临的挑战"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    content = slide.shapes.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    challenges = [
        ("⚖️ 监管挑战", [
            "• 平台规则不完善，治理手段存在漏洞",
            "• 监管滞后于技术创新速度",
            "• 数据安全与隐私保护问题"
        ]),
        ("💰 市场挑战", [
            "• 流量至上模式导致的恶性竞争",
            "• 大数据杀熟、算法歧视",
            "• 平台垄断与市场支配地位滥用"
        ]),
        ("🛒 运营挑战", [
            "• 刷单炒信、销售侵权商品",
            "• 低价竞争导致品质下降",
            "• 退货率上升影响消费体验"
        ]),
        ("🌐 发展挑战", [
            "• 增量市场开拓难度加大",
            "• 国际竞争加剧",
            "• 技术自主创新能力不足"
        ])
    ]
    
    for challenge_title, points in challenges:
        p = tf.add_paragraph()
        p.text = challenge_title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(204, 0, 0)  # 红色强调挑战
        p.space_after = Pt(6)
        
        for point in points:
            p = tf.add_paragraph()
            p.text = point
            p.font.size = Pt(18)
            p.level = 1
            p.space_before = Pt(3)
    
    # ========== 6. 高质量发展对策 ==========
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "四、平台经济高质量发展对策"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    content = slide.shapes.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    strategies = [
        ("🏛️ 完善监管体系", [
            "• 建立常态化监管制度",
            "• 制定差异化、精准化监管政策",
            "• 加强数据安全与隐私保护立法"
        ]),
        ("💡 鼓励技术创新", [
            "• 支持核心关键技术攻关",
            "• 促进平台企业加大研发投入",
            "• 建设开放创新的平台生态"
        ]),
        ("🤝 促进数实融合", [
            "• 推动平台经济赋能传统产业",
            "• 支持工业互联网平台发展",
            "• 建设数字化转型促进中心"
        ]),
        ("🌍 拓展国际市场", [
            "• 支持平台企业出海发展",
            "• 加强跨境电商平台建设",
            "• 参与国际规则制定"
        ]),
        ("⚖️ 优化营商环境", [
            "• 保障平台各类主体合法权益",
            "• 建立公平竞争的市场秩序",
            "• 完善平台治理机制"
        ])
    ]
    
    for strategy_title, points in strategies:
        p = tf.add_paragraph()
        p.text = strategy_title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 153, 76)  # 绿色强调对策
        p.space_after = Pt(6)
        
        for point in points:
            p = tf.add_paragraph()
            p.text = point
            p.font.size = Pt(18)
            p.level = 1
            p.space_before = Pt(3)
    
    # ========== 7. 未来展望 ==========
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "五、未来发展方向与展望"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    content = slide.shapes.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    future_outlook = [
        ("🔮 技术驱动", [
            "• AI、区块链、物联网等新技术深度应用",
            "• 平台智能化水平显著提升",
            "• 新型数字基础设施建设加速"
        ]),
        ("🔄 模式创新", [
            "• 平台商业模式持续创新",
            "• 共享经济、零工经济等新业态涌现",
            "• 线上线下深度融合"
        ]),
        ("🌐 全球化发展", [
            "• 平台企业国际化步伐加快",
            "• 数字贸易规则体系逐步完善",
            "• 全球数字治理合作加强"
        ]),
        ("🎯 高质量发展", [
            "• 从规模扩张转向质量提升",
            "• 更加注重社会责任与可持续发展",
            "• 成为经济高质量发展的重要支撑"
        ])
    ]
    
    for outlook_title, points in future_outlook:
        p = tf.add_paragraph()
        p.text = outlook_title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(102, 51, 204)  # 紫色
        p.space_after = Pt(6)
        
        for point in points:
            p = tf.add_paragraph()
            p.text = point
            p.font.size = Pt(18)
            p.level = 1
            p.space_before = Pt(3)
    
    # ========== 8. 结束页 ==========
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "谢谢！"
    subtitle.text = "Q&A"
    
    title.text_frame.paragraphs[0].font.size = Pt(48)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    subtitle.text_frame.paragraphs[0].font.size = Pt(32)
    subtitle.text_frame.paragraphs[0].font.italic = True
    
    # 保存PPT
    output_file = "平台经济高质量发展.pptx"
    prs.save(output_file)
    
    return output_file

if __name__ == "__main__":
    try:
        output_file = create_platform_economy_ppt()
        print(f"PPT创建成功：{output_file}")
        print(f"文件大小：{os.path.getsize(output_file) / 1024:.2f} KB")
    except Exception as e:
        print(f"创建PPT时出错：{e}")
        import traceback
        traceback.print_exc()