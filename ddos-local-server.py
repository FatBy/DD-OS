#!/usr/bin/env python3
"""
DD-OS Native Server v3.0
独立运行的本地 AI 操作系统后端

功能:
    - 文件操作 (读/写/列目录)
    - 命令执行 (Shell)
    - 任务管理 (后台执行)
    - 记忆持久化

用法:
    python ddos-local-server.py [--port 3001] [--path ~/clawd]

API:
    GET  /status              - 服务状态
    GET  /files               - 列出所有文件
    GET  /file/<name>         - 获取文件内容
    GET  /skills              - 获取技能列表
    GET  /memories            - 获取记忆数据
    GET  /all                 - 获取所有数据
    POST /api/tools/execute   - 执行工具 (新)
    POST /task/execute        - 执行任务 (兼容旧接口)
    GET  /task/status/<id>    - 查询任务状态
"""
from __future__ import annotations

import os
import sys
import re
import json
import argparse
import threading
import time
import uuid
import subprocess
import shlex
import shutil
from pathlib import Path
from http.server import ThreadingHTTPServer, BaseHTTPRequestHandler
from urllib.parse import unquote, urlparse, parse_qs
from datetime import datetime, timedelta
from concurrent.futures import ThreadPoolExecutor, Future

# PyYAML (skill-executor/parser.py 已依赖)
try:
    import yaml
    HAS_YAML = True
except ImportError:
    HAS_YAML = False

# MCP 客户端支持
try:
    from skills.mcp_manager import MCPClientManager
    HAS_MCP = True
except ImportError:
    HAS_MCP = False
    MCPClientManager = None

# 文件解析 (可选依赖，缺失时降级)
try:
    import pdfplumber
    HAS_PDF = True
except ImportError:
    HAS_PDF = False

try:
    from docx import Document as DocxDocument
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from pptx import Presentation as PptxPresentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    import pytesseract
    from PIL import Image
    HAS_OCR = True
except ImportError:
    HAS_OCR = False

import base64
import io

VERSION = "4.0.0"

# 🛡️ 安全配置
DANGEROUS_COMMANDS = {'rm -rf /', 'format', 'mkfs', 'dd if=/dev/zero'}
MAX_FILE_SIZE = 10 * 1024 * 1024  # 10MB 最大文件大小
MAX_OUTPUT_SIZE = 512 * 1024      # 512KB 最大输出
PLUGIN_TIMEOUT = 60               # 插件执行超时(秒)

# 🌐 静态文件 MIME 类型映射
MIME_TYPES = {
    '.html': 'text/html',
    '.css': 'text/css',
    '.js': 'application/javascript',
    '.json': 'application/json',
    '.png': 'image/png',
    '.jpg': 'image/jpeg',
    '.jpeg': 'image/jpeg',
    '.gif': 'image/gif',
    '.svg': 'image/svg+xml',
    '.ico': 'image/x-icon',
    '.woff': 'font/woff',
    '.woff2': 'font/woff2',
    '.ttf': 'font/ttf',
    '.eot': 'application/vnd.ms-fontobject',
}


# ============================================
# 🧮 文本相似度计算 (用于 Nexus 去重)
# ============================================

def calculate_text_similarity(text1: str, text2: str) -> float:
    """计算两个文本的 N-gram Jaccard 相似度"""
    if not text1 or not text2:
        return 0.0
    
    def get_ngrams(text: str) -> set:
        text = text.lower()
        # 清理符号，保留中英文和数字
        text = re.sub(r'[^\w\s\u4e00-\u9fff]', '', text)
        chars = list(text.replace(' ', ''))
        if len(chars) < 2:
            return set(chars)
        # 提取单字和相邻双字词 (Bi-gram)
        bigrams = [''.join(chars[i:i+2]) for i in range(len(chars)-1)]
        return set(chars + bigrams)
    
    set1 = get_ngrams(text1)
    set2 = get_ngrams(text2)
    if not set1 or not set2:
        return 0.0
    
    intersection = set1.intersection(set2)
    union = set1.union(set2)
    return len(intersection) / len(union)


# ============================================
# 🔌 SKILL.md Frontmatter 解析
# ============================================

def parse_skill_frontmatter(skill_md_path: Path) -> dict:
    """从 SKILL.md 提取 YAML frontmatter 元数据"""
    try:
        content = skill_md_path.read_text(encoding='utf-8')
    except Exception:
        return {}

    match = re.match(r'^---\s*\n(.*?)\n---\s*\n', content, re.DOTALL)
    if not match:
        # 无 frontmatter，提取第一段非标题文本作为 description
        desc = ''
        for line in content.split('\n'):
            line = line.strip()
            if line and not line.startswith('#'):
                desc = line[:200]
                break
        return {'description': desc}

    if HAS_YAML:
        try:
            return yaml.safe_load(match.group(1)) or {}
        except Exception:
            return {}
    else:
        # 无 PyYAML 时用简单正则提取 key: value
        result = {}
        for line in match.group(1).split('\n'):
            m = re.match(r'^(\w+)\s*:\s*(.+)$', line.strip())
            if m:
                key, val = m.group(1), m.group(2).strip()
                # 简单处理数组 [a, b, c]
                if val.startswith('[') and val.endswith(']'):
                    val = [v.strip().strip('"\'') for v in val[1:-1].split(',') if v.strip()]
                result[key] = val
        return result


def skill_name_to_tool_name(name: str) -> str:
    """将 skill 名称标准化为工具名 (kebab-case -> snake_case)"""
    return name.replace('-', '_').replace(' ', '_').lower()


# ============================================
# 🌌 NEXUS.md 解析
# ============================================

def parse_nexus_frontmatter(nexus_md_path: Path) -> dict:
    """从 NEXUS.md 提取 YAML frontmatter 元数据"""
    try:
        content = nexus_md_path.read_text(encoding='utf-8')
    except Exception:
        return {}

    match = re.match(r'^---\s*\n(.*?)\n---\s*\n', content, re.DOTALL)
    if not match:
        return {}

    if HAS_YAML:
        try:
            return yaml.safe_load(match.group(1)) or {}
        except Exception:
            return {}
    else:
        result = {}
        for line in match.group(1).split('\n'):
            m = re.match(r'^(\w+)\s*:\s*(.+)$', line.strip())
            if m:
                key, val = m.group(1), m.group(2).strip()
                if val.startswith('[') and val.endswith(']'):
                    val = [v.strip().strip('"\'') for v in val[1:-1].split(',') if v.strip()]
                result[key] = val
        return result


def extract_nexus_body(nexus_md_path: Path) -> str:
    """从 NEXUS.md 提取 Markdown 正文 (跳过 frontmatter)"""
    try:
        content = nexus_md_path.read_text(encoding='utf-8')
    except Exception:
        return ''

    # 去掉 frontmatter
    match = re.match(r'^---\s*\n.*?\n---\s*\n', content, re.DOTALL)
    if match:
        return content[match.end():].strip()
    return content.strip()


def update_nexus_frontmatter(nexus_md_path: Path, updates: dict):
    """更新 NEXUS.md 的 frontmatter 字段 (保留 body 不变)"""
    body = extract_nexus_body(nexus_md_path)
    frontmatter = parse_nexus_frontmatter(nexus_md_path)
    frontmatter.update(updates)

    # 重建 YAML frontmatter
    lines = ['---']
    for key, val in frontmatter.items():
        if isinstance(val, list):
            lines.append(f'{key}:')
            for item in val:
                lines.append(f'  - {item}')
        elif isinstance(val, dict):
            lines.append(f'{key}:')
            for k, v in val.items():
                lines.append(f'  {k}: {v}')
        else:
            lines.append(f'{key}: {val}')
    lines.append('---')
    lines.append('')
    lines.append(body)

    nexus_md_path.write_text('\n'.join(lines), encoding='utf-8')


def count_experience_entries(exp_dir: Path) -> int:
    """统计经验目录中的条目数，用于 XP 计算"""
    xp = 0
    successes = exp_dir / 'successes.md'
    failures = exp_dir / 'failures.md'
    if successes.exists():
        try:
            lines = successes.read_text(encoding='utf-8').split('\n')
            xp += sum(1 for l in lines if l.startswith('### ')) * 10
        except Exception:
            pass
    if failures.exists():
        try:
            lines = failures.read_text(encoding='utf-8').split('\n')
            xp += sum(1 for l in lines if l.startswith('### ')) * 5
        except Exception:
            pass
    return xp


# ============================================
# 🔌 动态工具注册表
# ============================================

class ToolRegistry:
    """动态工具发现与注册 - 支持内置工具 + 插件工具 + 指令型工具 + MCP工具"""

    def __init__(self, clawd_path: Path, project_path: Path = None):
        self.clawd_path = clawd_path
        # 项目目录 (脚本所在目录)，用于加载内置技能
        self.project_path = project_path or Path(__file__).parent.resolve()
        self.builtin_tools: dict = {}      # name -> callable
        self.plugin_tools: dict = {}       # name -> ToolSpec dict (有 execute.py)
        self.instruction_tools: dict = {}  # name -> InstructionSpec (纯 SKILL.md)
        self.mcp_tools: dict = {}          # name -> MCPToolSpec dict (MCP 服务器)
        self.mcp_manager: 'MCPClientManager | None' = None

    def register_builtin(self, name: str, handler):
        """注册内置工具"""
        self.builtin_tools[name] = handler

    def _get_skills_dirs(self) -> list[Path]:
        """获取所有技能目录 (用户目录 + 项目目录)"""
        dirs = []
        # 用户数据目录的技能 (优先级高，可覆盖内置)
        user_skills = self.clawd_path / 'skills'
        if user_skills.exists():
            dirs.append(user_skills)
        # 项目目录的内置技能
        project_skills = self.project_path / 'skills'
        if project_skills.exists() and project_skills != user_skills:
            dirs.append(project_skills)
        return dirs

    def scan_plugins(self):
        """递归扫描 skills/ 目录，注册可执行插件 + 指令型技能"""
        skills_dirs = self._get_skills_dirs()
        if not skills_dirs:
            return

        plugin_count = 0
        instruction_count = 0

        # 递归查找所有包含 SKILL.md 或 manifest.json 的目录
        seen_dirs: set = set()
        seen_tools: set = set()  # 防止重复注册同名工具

        for skills_dir in skills_dirs:
            # 1. 先扫描 manifest.json (可执行插件优先)
            for manifest_path in skills_dir.rglob('manifest.json'):
                skill_dir = manifest_path.parent
                dir_key = str(skill_dir.resolve())
                if dir_key in seen_dirs:
                    continue
                seen_dirs.add(dir_key)

                try:
                    spec = json.loads(manifest_path.read_text(encoding='utf-8'))

                    tools_list = spec.get('tools', [])
                    if not tools_list:
                        tools_list = [spec]

                    for tool_spec in tools_list:
                        tool_name = tool_spec.get('toolName', '')
                        executable = tool_spec.get('executable', spec.get('executable', 'execute.py'))

                        if not tool_name:
                            continue

                        # 跳过已注册的同名工具
                        if tool_name in seen_tools:
                            continue

                        exe_path = skill_dir / executable
                        if not exe_path.exists():
                            print(f"[ToolRegistry] Warning: {exe_path} not found, skipping {tool_name}")
                            continue

                        # 内置工具不可被覆盖
                        if tool_name in self.builtin_tools:
                            print(f"[ToolRegistry] Warning: plugin '{tool_name}' conflicts with builtin, skipping")
                            continue

                        self.plugin_tools[tool_name] = {
                            'name': tool_name,
                            'exe_path': str(exe_path),
                            'runtime': tool_spec.get('runtime', spec.get('runtime', 'python')),
                            'inputs': tool_spec.get('inputs', {}),
                            'outputs': tool_spec.get('outputs', {}),
                            'description': tool_spec.get('description', ''),
                            'dangerLevel': tool_spec.get('dangerLevel', spec.get('dangerLevel', 'safe')),
                            'version': tool_spec.get('version', spec.get('version', '1.0.0')),
                            'skill_dir': str(skill_dir),
                            'keywords': tool_spec.get('keywords', spec.get('keywords', [])),
                        }
                        seen_tools.add(tool_name)
                        plugin_count += 1
                        print(f"[ToolRegistry] Registered plugin: {tool_name} ({exe_path.name})")

                except Exception as e:
                    print(f"[ToolRegistry] Error loading {manifest_path}: {e}")

            # 2. 扫描 SKILL.md (指令型技能 - 没有 manifest.json 或没有 executable 的)
            for skill_md in skills_dir.rglob('SKILL.md'):
                skill_dir = skill_md.parent
                dir_key = str(skill_dir.resolve())

                # 已被 manifest.json 扫描注册的目录跳过
                if dir_key in seen_dirs:
                    continue
                seen_dirs.add(dir_key)

                try:
                    frontmatter = parse_skill_frontmatter(skill_md)
                    original_name = frontmatter.get('name', skill_dir.name)
                    tool_name = skill_name_to_tool_name(original_name)

                    # 冲突检查
                    if tool_name in self.builtin_tools or tool_name in self.plugin_tools or tool_name in seen_tools:
                        print(f"[ToolRegistry] Warning: instruction skill '{tool_name}' conflicts, skipping")
                        continue

                    self.instruction_tools[tool_name] = {
                        'name': tool_name,
                        'original_name': original_name,
                        'skill_path': str(skill_md),
                        'skill_dir': str(skill_dir),
                        'description': frontmatter.get('description', ''),
                        'inputs': frontmatter.get('inputs', {}),
                        'keywords': frontmatter.get('tags', frontmatter.get('keywords', [])),
                        'dangerLevel': 'safe',
                        'version': frontmatter.get('version', '1.0.0'),
                    }
                    seen_tools.add(tool_name)
                    instruction_count += 1
                    print(f"[ToolRegistry] Registered instruction skill: {tool_name} (from {skills_dir.name})")

                except Exception as e:
                    print(f"[ToolRegistry] Error loading {skill_md}: {e}")

        total = plugin_count + instruction_count
        if total > 0:
            print(f"[ToolRegistry] {total} tool(s) registered ({plugin_count} plugins, {instruction_count} instruction skills)")

    def scan_mcp_servers(self):
        """扫描并连接 MCP 服务器"""
        if not HAS_MCP:
            print("[ToolRegistry] MCP support not available (missing mcp_manager)")
            return

        self.mcp_manager = MCPClientManager(clawd_path=self.clawd_path)
        count = self.mcp_manager.initialize_all()

        if count > 0:
            # 注册 MCP 工具
            mcp_tool_count = 0
            for tool_info in self.mcp_manager.get_all_tools():
                tool_name = tool_info['name']
                # 冲突检查
                if tool_name in self.builtin_tools or tool_name in self.plugin_tools or tool_name in self.instruction_tools:
                    print(f"[ToolRegistry] Warning: MCP tool '{tool_name}' conflicts with existing tool, skipping")
                    continue

                self.mcp_tools[tool_name] = {
                    'name': tool_name,
                    'server': tool_info.get('server', ''),
                    'description': tool_info.get('description', ''),
                    'inputs': tool_info.get('inputs', {}),
                    'dangerLevel': 'safe',
                    'version': '1.0.0',
                }
                mcp_tool_count += 1

            print(f"[ToolRegistry] {mcp_tool_count} MCP tool(s) registered from {count} server(s)")

    def is_registered(self, name: str) -> bool:
        return name in self.builtin_tools or name in self.plugin_tools or name in self.instruction_tools or name in self.mcp_tools

    def get_plugin(self, name: str) -> dict | None:
        return self.plugin_tools.get(name)

    def get_instruction(self, name: str) -> dict | None:
        return self.instruction_tools.get(name)

    def get_mcp_tool(self, name: str) -> dict | None:
        return self.mcp_tools.get(name)

    def list_all(self) -> list:
        """返回所有已注册工具（内置+插件+指令型+MCP）"""
        # 内置工具元数据 (为有特殊参数的工具提供描述)
        BUILTIN_META = {
            'nexusBindSkill': {
                'description': '为当前 Nexus 绑定新技能依赖',
                'inputs': {
                    'nexusId': {'type': 'string', 'description': 'Nexus ID', 'required': True},
                    'skillId': {'type': 'string', 'description': '要绑定的技能 ID', 'required': True},
                },
            },
            'nexusUnbindSkill': {
                'description': '从当前 Nexus 移除技能依赖',
                'inputs': {
                    'nexusId': {'type': 'string', 'description': 'Nexus ID', 'required': True},
                    'skillId': {'type': 'string', 'description': '要移除的技能 ID', 'required': True},
                },
            },
            'parseFile': {
                'description': '解析文档文件（PDF/DOCX/PPTX）或对图像进行OCR文字识别，返回提取的文本内容',
                'inputs': {
                    'filePath': {'type': 'string', 'description': '文件路径（支持 .pdf .docx .pptx .png .jpg 等格式）', 'required': True},
                },
            },
            'generateSkill': {
                'description': '动态生成 Python SKILL 并保存。当现有工具无法完成任务时，用此工具创建新能力',
                'inputs': {
                    'name': {'type': 'string', 'description': '技能名称 (kebab-case，如 ppt-maker)', 'required': True},
                    'description': {'type': 'string', 'description': '技能功能描述', 'required': True},
                    'pythonCode': {'type': 'string', 'description': 'Python 实现代码（必须包含 main() 函数）', 'required': True},
                    'nexusId': {'type': 'string', 'description': '关联的 Nexus ID（可选，指定后保存到 Nexus 目录）', 'required': False},
                    'triggers': {'type': 'array', 'description': '触发关键词列表（可选）', 'required': False},
                },
            },
        }
        tools = []
        for name in self.builtin_tools:
            meta = BUILTIN_META.get(name, {})
            tools.append({'name': name, 'type': 'builtin', **meta})
        for name, spec in self.plugin_tools.items():
            tools.append({
                'name': name,
                'type': 'plugin',
                'description': spec.get('description', ''),
                'inputs': spec.get('inputs', {}),
                'dangerLevel': spec.get('dangerLevel', 'safe'),
                'version': spec.get('version', '1.0.0'),
            })
        for name, spec in self.instruction_tools.items():
            tools.append({
                'name': name,
                'type': 'instruction',
                'description': spec.get('description', ''),
                'inputs': spec.get('inputs', {}),
                'dangerLevel': 'safe',
                'version': spec.get('version', '1.0.0'),
            })
        for name, spec in self.mcp_tools.items():
            tools.append({
                'name': name,
                'type': 'mcp',
                'server': spec.get('server', ''),
                'description': spec.get('description', ''),
                'inputs': spec.get('inputs', {}),
                'dangerLevel': 'safe',
                'version': '1.0.0',
            })
        return tools


# ============================================
# 🤖 子代理管理器 (Quest 模式支持)
# ============================================

class SubagentManager:
    """
    子代理管理器 - 支持并行探索任务
    用于 Quest 模式的探索阶段，可同时运行多个轻量级代理
    """
    MAX_CONCURRENT = 5  # 最大并发数
    AGENT_TIMEOUT = 30  # 单个代理超时(秒)
    
    def __init__(self, tool_registry: ToolRegistry):
        self.registry = tool_registry
        self.agents: dict[str, dict] = {}  # agent_id -> agent_info
        self.executor = ThreadPoolExecutor(max_workers=self.MAX_CONCURRENT)
        self.futures: dict[str, Future] = {}  # agent_id -> Future
        self.lock = threading.Lock()
    
    def spawn(self, agent_type: str, task: str, tools: list[str], context: str = '') -> str:
        """
        启动一个子代理
        
        Args:
            agent_type: 代理类型 ('explore', 'plan', 'execute')
            task: 任务描述
            tools: 可用工具列表
            context: 上下文信息
        
        Returns:
            agent_id: 代理 ID
        """
        agent_id = f"subagent-{int(time.time()*1000)}-{uuid.uuid4().hex[:6]}"
        
        agent_info = {
            'id': agent_id,
            'type': agent_type,
            'task': task,
            'tools': tools,
            'context': context,
            'status': 'pending',
            'result': None,
            'error': None,
            'started_at': time.time(),
            'completed_at': None,
        }
        
        with self.lock:
            # 检查并发限制
            running_count = sum(1 for a in self.agents.values() if a['status'] == 'running')
            if running_count >= self.MAX_CONCURRENT:
                agent_info['status'] = 'queued'
                agent_info['error'] = f'Queue full, max {self.MAX_CONCURRENT} concurrent agents'
                self.agents[agent_id] = agent_info
                return agent_id
            
            self.agents[agent_id] = agent_info
        
        # 异步执行
        future = self.executor.submit(self._run_agent, agent_id, task, tools, context)
        self.futures[agent_id] = future
        
        with self.lock:
            self.agents[agent_id]['status'] = 'running'
        
        print(f"[SubagentManager] Spawned {agent_type} agent: {agent_id}")
        return agent_id
    
    def _run_agent(self, agent_id: str, task: str, tools: list[str], context: str) -> str:
        """
        执行子代理任务 (简化版 - 单工具调用)
        
        对于探索阶段，每个子代理通常只需要调用一个工具
        """
        try:
            result_parts = []
            
            # 根据任务类型选择工具
            for tool_name in tools:
                if not self.registry.is_registered(tool_name):
                    continue
                
                # 构建工具参数
                args = self._build_tool_args(tool_name, task, context)
                
                # 执行工具
                tool_result = self._execute_tool(tool_name, args)
                
                if tool_result.get('status') == 'success':
                    result_parts.append(f"[{tool_name}] {tool_result.get('result', '')[:1000]}")
                else:
                    result_parts.append(f"[{tool_name}] Error: {tool_result.get('result', 'Unknown error')[:200]}")
            
            final_result = '\n\n'.join(result_parts) if result_parts else 'No tools executed'
            
            with self.lock:
                if agent_id in self.agents:
                    self.agents[agent_id]['status'] = 'completed'
                    self.agents[agent_id]['result'] = final_result
                    self.agents[agent_id]['completed_at'] = time.time()
            
            return final_result
            
        except Exception as e:
            error_msg = str(e)
            with self.lock:
                if agent_id in self.agents:
                    self.agents[agent_id]['status'] = 'failed'
                    self.agents[agent_id]['error'] = error_msg
                    self.agents[agent_id]['completed_at'] = time.time()
            return f"Error: {error_msg}"
    
    def _build_tool_args(self, tool_name: str, task: str, context: str) -> dict:
        """根据工具类型构建参数"""
        # MCP quest 工具的特殊处理
        if tool_name == 'mcp__quest__search_codebase':
            # 提取关键词
            keywords = self._extract_keywords(task)
            return {
                'query': task,
                'key_words': ','.join(keywords[:3]),
                'explanation': f'Exploring: {task[:50]}'
            }
        elif tool_name == 'mcp__quest__search_symbol':
            # 从任务中提取符号名
            symbols = self._extract_symbols(task)
            return {
                'queries': [{'symbol': s, 'relation': 'all'} for s in symbols[:2]],
                'explanation': f'Symbol search for: {task[:50]}'
            }
        elif tool_name == 'readFile':
            # 从上下文中提取文件路径
            paths = self._extract_file_paths(context)
            return {'path': paths[0] if paths else ''}
        elif tool_name == 'listDir':
            return {'path': '.', 'recursive': False}
        else:
            return {'query': task}
    
    def _extract_keywords(self, text: str) -> list[str]:
        """从文本中提取关键词"""
        # 简单实现：提取英文单词和中文词组
        words = re.findall(r'[a-zA-Z_][a-zA-Z0-9_]*|[\u4e00-\u9fff]+', text)
        # 过滤常见词
        stopwords = {'the', 'a', 'an', 'is', 'are', 'to', 'for', 'of', 'in', 'on', 'with', '的', '是', '在', '和', '了'}
        return [w for w in words if w.lower() not in stopwords and len(w) > 1][:5]
    
    def _extract_symbols(self, text: str) -> list[str]:
        """从文本中提取可能的符号名（函数名、类名等）"""
        # 匹配驼峰命名和下划线命名
        symbols = re.findall(r'\b([A-Z][a-zA-Z0-9]*|[a-z][a-zA-Z0-9]*(?:_[a-zA-Z0-9]+)+)\b', text)
        return list(set(symbols))[:3]
    
    def _extract_file_paths(self, text: str) -> list[str]:
        """从文本中提取文件路径"""
        paths = re.findall(r'[a-zA-Z0-9_./\\-]+\.[a-zA-Z]+', text)
        return paths[:3]
    
    def _execute_tool(self, tool_name: str, args: dict) -> dict:
        """执行单个工具"""
        # 检查 MCP 工具
        if tool_name.startswith('mcp__') and self.registry.mcp_manager:
            try:
                result = self.registry.mcp_manager.call_tool(tool_name, args)
                return {'status': 'success', 'result': str(result)[:2000]}
            except Exception as e:
                return {'status': 'error', 'result': str(e)}
        
        # 内置工具需要通过 handler 执行，这里返回占位
        return {'status': 'error', 'result': f'Tool {tool_name} not directly executable in subagent'}
    
    def get_status(self, agent_id: str) -> dict | None:
        """获取子代理状态"""
        with self.lock:
            return self.agents.get(agent_id)
    
    def get_all_status(self) -> list[dict]:
        """获取所有子代理状态"""
        with self.lock:
            return list(self.agents.values())
    
    def collect_results(self, agent_ids: list[str], timeout: float = 60.0) -> list[dict]:
        """
        收集多个子代理的结果
        
        Args:
            agent_ids: 要收集的代理 ID 列表
            timeout: 等待超时时间(秒)
        
        Returns:
            结果列表
        """
        results = []
        deadline = time.time() + timeout
        
        for agent_id in agent_ids:
            remaining = deadline - time.time()
            if remaining <= 0:
                break
            
            future = self.futures.get(agent_id)
            if future:
                try:
                    future.result(timeout=remaining)
                except Exception:
                    pass
            
            with self.lock:
                agent = self.agents.get(agent_id)
                if agent:
                    results.append({
                        'id': agent_id,
                        'type': agent['type'],
                        'task': agent['task'],
                        'status': agent['status'],
                        'result': agent.get('result'),
                        'error': agent.get('error'),
                    })
        
        return results
    
    def cleanup_old_agents(self, max_age: float = 300.0):
        """清理超过指定时间的旧代理记录"""
        cutoff = time.time() - max_age
        with self.lock:
            to_remove = [
                aid for aid, agent in self.agents.items()
                if agent.get('completed_at', 0) < cutoff and agent['status'] in ('completed', 'failed')
            ]
            for aid in to_remove:
                del self.agents[aid]
                self.futures.pop(aid, None)
        
        if to_remove:
            print(f"[SubagentManager] Cleaned up {len(to_remove)} old agents")


class ClawdDataHandler(BaseHTTPRequestHandler):
    clawd_path = None
    project_path = None  # 项目目录，用于加载内置技能
    registry = None  # type: ToolRegistry
    subagent_manager = None  # type: SubagentManager
    tasks = {}
    tasks_lock = threading.Lock()
    
    def log_message(self, format, *args):
        timestamp = datetime.now().strftime('%H:%M:%S')
        print(f"[{timestamp}] {format % args}")
    
    def send_cors_headers(self):
        self.send_header('Access-Control-Allow-Origin', '*')
        self.send_header('Access-Control-Allow-Methods', 'GET, POST, OPTIONS')
        self.send_header('Access-Control-Allow-Headers', 'Content-Type')
    
    def send_json(self, data, status=200):
        self.send_response(status)
        self.send_header('Content-Type', 'application/json')
        self.send_cors_headers()
        self.end_headers()
        self.wfile.write(json.dumps(data, ensure_ascii=False, indent=2).encode('utf-8'))
    
    def send_text(self, text, status=200):
        self.send_response(status)
        self.send_header('Content-Type', 'text/plain; charset=utf-8')
        self.send_cors_headers()
        self.end_headers()
        self.wfile.write(text.encode('utf-8'))
    
    def send_error_json(self, message, status=404):
        self.send_json({'error': message, 'status': 'error'}, status)
    
    def do_OPTIONS(self):
        self.send_response(200)
        self.send_cors_headers()
        self.end_headers()
    
    def do_GET(self):
        parsed = urlparse(self.path)
        path = unquote(parsed.path)
        query = parse_qs(parsed.query)
        
        routes = {
            '/status': self.handle_status,
            '/files': self.handle_files,
            '/skills': self.handle_skills,
            '/nexuses': self.handle_nexuses,
            '/memories': self.handle_memories,
            '/tools': self.handle_tools_list,
            '/all': self.handle_all,
            '/': self.handle_index,
            '': self.handle_index,
        }
        
        if path in routes:
            routes[path]()
        elif path.startswith('/file/'):
            self.handle_file(path[6:])
        elif path.startswith('/nexuses/') and '/experience' not in path:
            nexus_name = path[9:]  # strip '/nexuses/'
            if nexus_name == 'health':
                self.handle_nexuses_health()
            else:
                self.handle_nexus_detail(nexus_name)
        elif path.startswith('/task/status/'):
            task_id = path[13:]
            offset = int(query.get('offset', ['0'])[0])
            self.handle_task_status(task_id, offset)
        elif path == '/api/traces/search':
            self.handle_trace_search(query)
        elif path == '/api/traces/recent':
            self.handle_trace_recent(query)
        elif path == '/api/registry/skills':
            self.handle_registry_skills_search(query)
        elif path == '/api/registry/mcp':
            self.handle_registry_mcp_search(query)
        elif path == '/mcp/servers':
            self.handle_mcp_servers_list()
        elif path.startswith('/data/'):
            # 前端数据读取 API
            key = path[6:]  # strip '/data/'
            self.handle_data_get(key)
        elif path == '/data':
            # 列出所有数据键
            self.handle_data_list()
        else:
            # 静态文件服务 (托管 dist/ 目录)
            self.serve_static_file(path)
    
    def do_POST(self):
        parsed = urlparse(self.path)
        path = unquote(parsed.path)
        content_type = self.headers.get('Content-Type', '')
        
        # 文件上传：multipart/form-data 单独处理（避免大文件 JSON 编码 OOM）
        if path == '/api/files/upload' and 'multipart/form-data' in content_type:
            self.handle_file_upload_multipart()
            return
        
        content_length = int(self.headers.get('Content-Length', 0))
        body = self.rfile.read(content_length).decode('utf-8') if content_length > 0 else '{}'
        
        try:
            data = json.loads(body) if body else {}
        except json.JSONDecodeError:
            self.send_error_json('Invalid JSON', 400)
            return
        
        # 🌟 新增：工具执行接口
        if path == '/api/tools/execute':
            self.handle_tool_execution(data)
        elif path == '/api/files/upload':
            self.handle_file_upload(data)
        elif path == '/tools/reload':
            self.handle_tools_reload(data)
        elif path == '/api/traces/save':
            self.handle_trace_save(data)
        elif path == '/mcp/reload':
            self.handle_mcp_reload(data)
        elif path.startswith('/mcp/servers/') and path.endswith('/reconnect'):
            server_name = path[13:-10]  # Extract server name
            self.handle_mcp_reconnect(server_name)
        elif path == '/mcp/install':
            self.handle_mcp_install(data)
        elif path == '/skills/install':
            self.handle_skill_install(data)
        elif path == '/skills/uninstall':
            self.handle_skill_uninstall(data)
        elif path.startswith('/nexuses/') and path.endswith('/skills'):
            nexus_name = path[9:-7]  # strip '/nexuses/' and '/skills'
            self.handle_nexus_update_skills(nexus_name, data)
        elif path.startswith('/nexuses/') and path.endswith('/experience'):
            nexus_name = path[9:-11]  # strip '/nexuses/' and '/experience'
            self.handle_add_experience(nexus_name, data)
        elif path.startswith('/nexuses/') and path.endswith('/meta'):
            nexus_name = path[9:-5]  # strip '/nexuses/' and '/meta'
            self.handle_nexus_update_meta(nexus_name, data)
        elif path == '/task/execute':
            self.handle_task_execute(data)
        elif path.startswith('/data/'):
            # 前端数据写入 API
            key = path[6:]  # strip '/data/'
            self.handle_data_set(key, data)
        # 🤖 子代理 API (Quest 模式支持)
        elif path == '/api/subagent/spawn':
            self.handle_subagent_spawn(data)
        elif path == '/api/subagent/collect':
            self.handle_subagent_collect(data)
        elif path.startswith('/api/subagent/') and path.endswith('/status'):
            agent_id = path[14:-7]  # strip '/api/subagent/' and '/status'
            self.handle_subagent_status(agent_id)
        else:
            self.send_error_json(f'Unknown endpoint: {path}', 404)
    
    # ============================================
    # 🌐 静态文件服务 (托管前端 dist/)
    # ============================================
    
    def serve_static_file(self, path: str):
        """托管 dist/ 目录的前端构建产物，支持 SPA 路由"""
        # 静态文件目录 (与服务器脚本同级的 dist/)
        static_dir = Path(__file__).parent / 'dist'
        
        if not static_dir.exists():
            # dist/ 不存在时返回提示
            self.send_response(200)
            self.send_header('Content-Type', 'text/html; charset=utf-8')
            self.end_headers()
            self.wfile.write(b'''<!DOCTYPE html>
<html>
<head><title>DD-OS Server</title></head>
<body style="font-family: system-ui; padding: 40px; background: #1a1a2e; color: #eee;">
<h1>DD-OS Native Server</h1>
<p>Frontend not built. Run <code>npm run build</code> to generate dist/</p>
<p>Or access dev server at <a href="http://localhost:5173">http://localhost:5173</a></p>
<hr>
<p>API Endpoints:</p>
<ul>
<li>GET /status - Server status</li>
<li>GET /skills - List skills</li>
<li>POST /api/tools/execute - Execute tool</li>
</ul>
</body>
</html>''')
            return
        
        # 确定文件路径
        if path == '/' or path == '':
            file_path = static_dir / 'index.html'
        else:
            # 去掉开头的 /
            clean_path = path.lstrip('/')
            file_path = static_dir / clean_path
        
        # SPA 路由支持：如果不是文件（没有扩展名），返回 index.html
        if not file_path.exists():
            if '.' not in file_path.name:
                file_path = static_dir / 'index.html'
        
        if not file_path.exists():
            self.send_error_json(f'File not found: {path}', 404)
            return
        
        # 安全检查：确保路径在 static_dir 内
        try:
            file_path.resolve().relative_to(static_dir.resolve())
        except ValueError:
            self.send_error_json('Access denied', 403)
            return
        
        # 获取 MIME 类型
        suffix = file_path.suffix.lower()
        content_type = MIME_TYPES.get(suffix, 'application/octet-stream')
        
        try:
            with open(file_path, 'rb') as f:
                content = f.read()
            
            self.send_response(200)
            self.send_header('Content-Type', content_type)
            self.send_header('Content-Length', str(len(content)))
            # 缓存控制：静态资源长期缓存
            if '/assets/' in str(file_path):
                self.send_header('Cache-Control', 'public, max-age=31536000')
            self.end_headers()
            self.wfile.write(content)
        except Exception as e:
            self.send_error_json(f'Failed to read file: {str(e)}', 500)
    
    # ============================================
    # 📦 前端数据持久化 API (/data)
    # ============================================
    
    def handle_data_get(self, key: str):
        """读取前端数据"""
        data_dir = self.clawd_path / 'data'
        data_dir.mkdir(exist_ok=True)
        
        # 安全检查：key 只能是字母数字下划线
        if not re.match(r'^[a-zA-Z0-9_-]+$', key):
            self.send_error_json('Invalid key format', 400)
            return
        
        file_path = data_dir / f'{key}.json'
        
        if not file_path.exists():
            self.send_json({'key': key, 'value': None, 'exists': False})
            return
        
        try:
            content = file_path.read_text(encoding='utf-8')
            self.send_json({'key': key, 'value': json.loads(content), 'exists': True})
        except Exception as e:
            self.send_error_json(f'Failed to read data: {str(e)}', 500)
    
    def handle_data_set(self, key: str, data: dict):
        """写入前端数据"""
        data_dir = self.clawd_path / 'data'
        data_dir.mkdir(exist_ok=True)
        
        # 安全检查：key 只能是字母数字下划线
        if not re.match(r'^[a-zA-Z0-9_-]+$', key):
            self.send_error_json('Invalid key format', 400)
            return
        
        file_path = data_dir / f'{key}.json'
        value = data.get('value')
        
        try:
            if value is None:
                # 删除数据
                if file_path.exists():
                    file_path.unlink()
                self.send_json({'key': key, 'deleted': True})
            else:
                # 写入数据
                file_path.write_text(json.dumps(value, ensure_ascii=False, indent=2), encoding='utf-8')
                self.send_json({'key': key, 'saved': True})
        except Exception as e:
            self.send_error_json(f'Failed to save data: {str(e)}', 500)
    
    def handle_data_list(self):
        """列出所有数据键"""
        data_dir = self.clawd_path / 'data'
        data_dir.mkdir(exist_ok=True)
        
        keys = []
        for f in data_dir.glob('*.json'):
            keys.append(f.stem)
        
        self.send_json({'keys': keys})
    
    # ============================================
    # 🛠️ 工具执行 (核心新功能)
    # ============================================
    
    def handle_tool_execution(self, data):
        """处理工具调用请求 - 支持内置工具、插件工具、指令型工具和MCP工具"""
        tool_name = data.get('name', '')
        args = data.get('args', {})

        if not self.registry.is_registered(tool_name):
            all_tools = [t['name'] for t in self.registry.list_all()]
            self.send_json({
                'tool': tool_name,
                'status': 'error',
                'result': f'Tool not registered: {tool_name}. Available: {", ".join(all_tools)}'
            }, 403)
            return

        result = ""
        status = "success"

        try:
            # 1. 指令型工具 -> 路由到 skill-executor
            instruction_spec = self.registry.get_instruction(tool_name)
            if instruction_spec:
                result = self._execute_instruction_tool(instruction_spec, tool_name, args)
            # 2. 插件工具 -> subprocess 执行
            elif self.registry.get_plugin(tool_name):
                plugin_spec = self.registry.get_plugin(tool_name)
                result = self._execute_plugin_tool(plugin_spec, tool_name, args)
            # 3. MCP 工具 -> 通过 MCPManager 调用
            elif self.registry.get_mcp_tool(tool_name):
                result = self._execute_mcp_tool(tool_name, args)
            # 4. 内置工具 -> 直接调度
            else:
                builtin_handlers = {
                    'readFile': self._tool_read_file,
                    'writeFile': self._tool_write_file,
                    'appendFile': self._tool_append_file,
                    'listDir': self._tool_list_dir,
                    'runCmd': self._tool_run_cmd,
                    'weather': self._tool_weather,
                    'webSearch': self._tool_web_search,
                    'webFetch': self._tool_web_fetch,
                    'saveMemory': self._tool_save_memory,
                    'searchMemory': self._tool_search_memory,
                    'nexusBindSkill': self._tool_nexus_bind_skill,
                    'nexusUnbindSkill': self._tool_nexus_unbind_skill,
                    'openInExplorer': self._tool_open_in_explorer,
                    'parseFile': self._tool_parse_file,
                    'generateSkill': self._tool_generate_skill,
                }
                handler = builtin_handlers.get(tool_name)
                if handler:
                    result = handler(args)
                else:
                    raise ValueError(f"No handler for builtin tool: {tool_name}")

        except Exception as e:
            status = "error"
            result = f"Tool execution failed: {str(e)}"

        self.send_json({
            'tool': tool_name,
            'status': status,
            'result': result,
            'timestamp': datetime.now().isoformat()
        })

    # ============================================
    # 📎 文件上传 + 自动解析
    # ============================================

    UPLOAD_ALLOWED_EXT = {'.pdf', '.docx', '.pptx', '.png', '.jpg', '.jpeg', '.bmp', '.tiff', '.webp', '.txt', '.md', '.csv'}

    def handle_file_upload_multipart(self):
        """接收 FormData multipart 上传，保存并自动解析
        
        使用手动 multipart boundary 解析，不依赖已废弃的 cgi.FieldStorage
        （cgi 在 Python 3.11+ deprecated，3.13 removed）
        """
        content_type = self.headers.get('Content-Type', '')
        content_length = int(self.headers.get('Content-Length', 0))

        if content_length > MAX_FILE_SIZE:
            # 消耗请求体避免连接异常
            remaining = content_length
            while remaining > 0:
                chunk = min(remaining, 65536)
                self.rfile.read(chunk)
                remaining -= chunk
            self.send_error_json(f'文件过大 (>{MAX_FILE_SIZE // 1024 // 1024}MB)', 413)
            return

        # 从 Content-Type 提取 boundary
        boundary = None
        for part in content_type.split(';'):
            part = part.strip()
            if part.startswith('boundary='):
                boundary = part[len('boundary='):].strip('"')
        if not boundary:
            self.send_error_json('无效的 multipart 请求：缺少 boundary', 400)
            return

        # 读取整个请求体（已验证 content_length <= 10MB，内存安全）
        try:
            raw_body = self.rfile.read(content_length)
        except Exception as e:
            print(f"[ERROR] 读取请求体失败: {e}", file=sys.stderr)
            self.send_error_json('读取上传数据失败', 400)
            return

        # 按 boundary 分割，提取包含 filename 的 part
        boundary_bytes = ('--' + boundary).encode()
        parts = raw_body.split(boundary_bytes)

        file_bytes = None
        file_name = 'unknown'
        for part_data in parts:
            if b'filename=' not in part_data:
                continue
            # headers 和 body 以空行 (\r\n\r\n) 分隔
            header_end = part_data.find(b'\r\n\r\n')
            if header_end == -1:
                continue
            headers_raw_bytes = part_data[:header_end]
            # 尝试 UTF-8（浏览器 FormData），回退到 GBK（Windows curl/工具）
            try:
                headers_raw = headers_raw_bytes.decode('utf-8')
            except UnicodeDecodeError:
                headers_raw = headers_raw_bytes.decode('gbk', errors='replace')
            file_bytes = part_data[header_end + 4:]
            # 去掉尾部的 \r\n（multipart 格式约定）
            if file_bytes.endswith(b'\r\n'):
                file_bytes = file_bytes[:-2]
            # 从 Content-Disposition 提取 filename
            for line in headers_raw.split('\r\n'):
                if 'filename=' in line:
                    # 支持: filename="中文.pptx" 和 filename=file.pdf
                    match = re.search(r'filename="?([^";\r\n]+)"?', line)
                    if match:
                        file_name = match.group(1).strip()
            break  # 只取第一个文件

        if file_bytes is None:
            self.send_error_json('未找到上传文件', 400)
            return

        # 清理文件名（保留中文、字母、数字、点、横线）
        safe_name = re.sub(r'[^\w.\-\u4e00-\u9fff]', '_', file_name)
        # 文件名长度限制（NTFS/ext4 最大 255 字符）
        stem, ext = os.path.splitext(safe_name)
        ext = ext.lower()
        if len(safe_name) > 200:
            safe_name = stem[:200 - len(ext)] + ext

        if ext not in self.UPLOAD_ALLOWED_EXT:
            self.send_error_json(f'不支持的文件类型: {ext}，支持: {", ".join(sorted(self.UPLOAD_ALLOWED_EXT))}', 400)
            return

        if len(file_bytes) > MAX_FILE_SIZE:
            self.send_error_json(f'文件过大 (>{MAX_FILE_SIZE // 1024 // 1024}MB)', 413)
            return

        # 保存到临时目录
        upload_dir = self.clawd_path / 'temp' / 'uploads'
        upload_dir.mkdir(parents=True, exist_ok=True)
        unique_name = f"{uuid.uuid4().hex[:8]}_{safe_name}"
        file_path = upload_dir / unique_name

        try:
            file_path.write_bytes(file_bytes)
        except Exception as e:
            print(f"[ERROR] 文件保存失败: {e}", file=sys.stderr)
            self.send_error_json('文件保存失败', 500)
            return

        # 自动解析
        parsed_text = ''
        try:
            parsed_text = self._tool_parse_file({'filePath': str(file_path)})
        except Exception as e:
            print(f"[ERROR] 文件解析失败: {e}", file=sys.stderr)
            parsed_text = f'[解析失败: 请检查文件格式是否正确]'

        file_size = len(file_bytes)
        self.send_json({
            'success': True,
            'filePath': str(file_path),
            'originalName': file_name,
            'fileSize': file_size,
            'parsedText': parsed_text,
            'timestamp': datetime.now().isoformat()
        })

    def handle_file_upload(self, data: dict):
        """接收前端上传的文件（Base64），保存到临时目录并自动解析"""
        file_name = data.get('fileName', '')
        data_base64 = data.get('dataBase64', '')

        if not file_name or not data_base64:
            self.send_error_json('fileName and dataBase64 are required', 400)
            return

        # 清理文件名
        safe_name = re.sub(r'[^\w.\-\u4e00-\u9fff]', '_', file_name)
        ext = os.path.splitext(safe_name)[1].lower()

        if ext not in self.UPLOAD_ALLOWED_EXT:
            self.send_error_json(f'不支持的文件类型: {ext}，支持: {", ".join(sorted(self.UPLOAD_ALLOWED_EXT))}', 400)
            return

        # 解码 Base64 (去掉 data:xxx;base64, 前缀)
        try:
            if ';base64,' in data_base64:
                data_base64 = data_base64.split(';base64,')[1]
            file_bytes = base64.b64decode(data_base64)
        except Exception as e:
            self.send_error_json(f'Base64 解码失败: {str(e)}', 400)
            return

        if len(file_bytes) > MAX_FILE_SIZE:
            self.send_error_json(f'文件过大 (>{MAX_FILE_SIZE // 1024 // 1024}MB)', 413)
            return

        # 保存到临时目录
        upload_dir = self.clawd_path / 'temp' / 'uploads'
        upload_dir.mkdir(parents=True, exist_ok=True)
        unique_name = f"{uuid.uuid4().hex[:8]}_{safe_name}"
        file_path = upload_dir / unique_name

        try:
            file_path.write_bytes(file_bytes)
        except Exception as e:
            self.send_error_json(f'文件保存失败: {str(e)}', 500)
            return

        # 自动解析
        parsed_text = ''
        try:
            parsed_text = self._tool_parse_file({'filePath': str(file_path)})
        except Exception as e:
            parsed_text = f'[解析失败: {str(e)}]'

        self.send_json({
            'success': True,
            'filePath': str(file_path),
            'originalName': file_name,
            'parsedText': parsed_text,
            'timestamp': datetime.now().isoformat()
        })

    def _execute_plugin_tool(self, spec: dict, tool_name: str, args: dict) -> str:
        """执行插件工具 - subprocess 隔离执行"""
        exe_path = spec['exe_path']
        runtime = spec.get('runtime', 'python')

        # 确定运行时命令
        if runtime == 'python':
            cmd = [sys.executable, exe_path]
        elif runtime == 'node':
            cmd = ['node', exe_path]
        else:
            raise ValueError(f"Unsupported runtime: {runtime}")

        # 构建输入：包含工具名和参数（支持多工具 manifest）
        input_data = json.dumps({
            'tool': tool_name,
            'args': args
        }, ensure_ascii=False)

        try:
            process = subprocess.run(
                cmd,
                input=input_data,
                capture_output=True,
                text=True,
                timeout=PLUGIN_TIMEOUT,
                cwd=spec.get('skill_dir', str(self.clawd_path)),
            )

            if process.returncode != 0:
                stderr = process.stderr[:MAX_OUTPUT_SIZE] if process.stderr else ''
                raise RuntimeError(f"Plugin exited with code {process.returncode}: {stderr}")

            return process.stdout[:MAX_OUTPUT_SIZE] if process.stdout else ''

        except subprocess.TimeoutExpired:
            raise RuntimeError(f"Plugin timed out after {PLUGIN_TIMEOUT}s")

    def _execute_instruction_tool(self, spec: dict, tool_name: str, args: dict) -> str:
        """执行指令型工具 - 通过 skill-executor 解析 SKILL.md 并返回指令"""
        skill_executor = self.clawd_path / 'skills' / 'skill-executor' / 'execute.py'

        if not skill_executor.exists():
            raise RuntimeError(f"skill-executor not found at {skill_executor}")

        # 使用 original_name (kebab-case) 让 SkillDiscovery 能找到目录
        original_name = spec.get('original_name', tool_name)

        input_data = json.dumps({
            'tool': 'run_skill',
            'args': {
                'skill_name': original_name,
                'args': args,
                'project_root': str(self.clawd_path),
            }
        }, ensure_ascii=False)

        try:
            process = subprocess.run(
                [sys.executable, str(skill_executor)],
                input=input_data,
                capture_output=True,
                text=True,
                timeout=PLUGIN_TIMEOUT,
                cwd=str(skill_executor.parent),
            )

            if process.returncode != 0:
                stderr = process.stderr[:MAX_OUTPUT_SIZE] if process.stderr else ''
                raise RuntimeError(f"Instruction skill error: {stderr}")

            result = json.loads(process.stdout)
            if not result.get('success'):
                raise RuntimeError(result.get('error', 'Unknown error'))

            return result.get('instructions', result.get('output', ''))

        except subprocess.TimeoutExpired:
            raise RuntimeError(f"Instruction skill timed out after {PLUGIN_TIMEOUT}s")
        except json.JSONDecodeError:
            # skill-executor 返回非 JSON 时，直接返回原文
            return process.stdout[:MAX_OUTPUT_SIZE] if process.stdout else ''

    def _execute_mcp_tool(self, tool_name: str, args: dict) -> str:
        """执行 MCP 工具 - 通过 MCPManager 调用远程 MCP 服务器"""
        if not self.registry.mcp_manager:
            raise RuntimeError("MCP manager not initialized")

        try:
            result = self.registry.mcp_manager.call_tool(tool_name, args, timeout=PLUGIN_TIMEOUT)
            if result is None:
                return ""
            return str(result)
        except Exception as e:
            raise RuntimeError(f"MCP tool execution failed: {e}")
    
    def _resolve_path(self, relative_path: str, allow_outside: bool = False) -> Path:
        """解析并验证路径安全性"""
        if not relative_path:
            raise ValueError("Path cannot be empty")
        
        # 移除开头的斜杠
        clean_path = relative_path.lstrip('/')
        
        # 默认在 clawd 目录下操作
        if allow_outside and os.path.isabs(relative_path):
            file_path = Path(relative_path)
        else:
            file_path = self.clawd_path / clean_path
        
        # 安全检查：防止路径遍历
        try:
            resolved = file_path.resolve()
            if not allow_outside:
                resolved.relative_to(self.clawd_path.resolve())
        except ValueError:
            raise PermissionError(f"Access denied: path outside allowed directory")
        
        return resolved
    
    def _tool_read_file(self, args: dict) -> str:
        """读取文件内容"""
        path = args.get('path', '')
        file_path = self._resolve_path(path, allow_outside=args.get('allowOutside', False))
        
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {path}")
        if not file_path.is_file():
            raise ValueError(f"Not a file: {path}")
        if file_path.stat().st_size > MAX_FILE_SIZE:
            raise ValueError(f"File too large (>{MAX_FILE_SIZE} bytes)")
        
        return file_path.read_text(encoding='utf-8')
    
    def _tool_parse_file(self, args: dict) -> str:
        """解析文档或图像文件，返回提取的文本内容"""
        file_path_str = args.get('filePath') or args.get('path', '')
        if not file_path_str:
            raise ValueError("filePath is required")
        
        # 支持绝对路径（上传的临时文件）和相对路径
        if os.path.isabs(file_path_str):
            file_path = Path(file_path_str).resolve()
            # 安全检查：只允许访问 clawd 工作目录下的文件
            allowed_root = self.clawd_path.resolve()
            try:
                file_path.relative_to(allowed_root)
            except ValueError:
                raise PermissionError(f"Access denied: path outside allowed directory")
        else:
            file_path = self._resolve_path(file_path_str, allow_outside=True)
        
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path_str}")
        if file_path.stat().st_size > MAX_FILE_SIZE:
            raise ValueError(f"File too large (>{MAX_FILE_SIZE // 1024 // 1024}MB)")
        
        ext = file_path.suffix.lower()
        text = ""
        
        if ext == '.pdf':
            if not HAS_PDF:
                raise RuntimeError("pdfplumber 未安装，请运行 pip install pdfplumber")
            with pdfplumber.open(str(file_path)) as pdf:
                pages = []
                for i, page in enumerate(pdf.pages):
                    page_text = page.extract_text() or ''
                    if page_text.strip():
                        pages.append(f"--- 第{i+1}页 ---\n{page_text}")
                text = "\n\n".join(pages)
        
        elif ext == '.docx':
            if not HAS_DOCX:
                raise RuntimeError("python-docx 未安装，请运行 pip install python-docx")
            doc = DocxDocument(str(file_path))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            for table in doc.tables:
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    paragraphs.append(" | ".join(cells))
            text = "\n".join(paragraphs)
        
        elif ext == '.pptx':
            if not HAS_PPTX:
                raise RuntimeError("python-pptx 未安装，请运行 pip install python-pptx")
            prs = PptxPresentation(str(file_path))
            slides = []
            for i, slide in enumerate(prs.slides):
                parts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                parts.append(t)
                if parts:
                    slides.append(f"--- 幻灯片{i+1} ---\n" + "\n".join(parts))
            text = "\n\n".join(slides)
        
        elif ext in ('.png', '.jpg', '.jpeg', '.bmp', '.tiff', '.webp'):
            if not HAS_OCR:
                raise RuntimeError("pytesseract/Pillow 未安装，请运行 pip install pytesseract Pillow")
            img = Image.open(str(file_path))
            lang = args.get('language', 'eng+chi_sim')
            text = pytesseract.image_to_string(img, lang=lang)
        
        else:
            # 回退：尝试当纯文本读取
            try:
                text = file_path.read_text(encoding='utf-8')
            except Exception:
                raise ValueError(f"不支持的文件类型: {ext}")
        
        if not text.strip():
            return f"[文件 {file_path.name} 无可提取的文本内容]"
        
        # 截断到 MAX_OUTPUT_SIZE（安全 UTF-8 边界截断）
        encoded = text.encode('utf-8')
        if len(encoded) > MAX_OUTPUT_SIZE:
            safe_idx = MAX_OUTPUT_SIZE
            # 回退到 UTF-8 字符边界，避免截断多字节字符导致乱码
            while safe_idx > 0 and (encoded[safe_idx] & 0xC0) == 0x80:
                safe_idx -= 1
            text = encoded[:safe_idx].decode('utf-8')
            text += f"\n\n[内容过长，已截断至约 {MAX_OUTPUT_SIZE // 1024}KB]"
        
        return text
    
    def _tool_write_file(self, args: dict) -> str:
        """写入文件"""
        path = args.get('path', '')
        content = args.get('content', '')
        
        file_path = self._resolve_path(path)
        
        # === Nexus 涌现去重网关 ===
        if 'nexuses/' in path and path.endswith('NEXUS.md'):
            # 仅在文件不存在时（即新建操作）进行去重检查
            if not file_path.exists():
                duplicate_id = self._check_nexus_duplication(content)
                if duplicate_id:
                    return (f"【系统拦截】创建失败！\n"
                            f"检测到高度相似的 Nexus 节点已存在 (节点 ID: {duplicate_id})。\n"
                            f"为避免知识图谱碎片化，请不要创建新目录，请直接使用 'readFile' 和 'writeFile' "
                            f"读取并更新原有的 nexuses/{duplicate_id}/NEXUS.md，或者向其追加 experience。")
        
        # === Nexus 格式引导 ===
        # 检测写入 nexuses/ 目录但不是 NEXUS.md 的情况，提供格式纠正提示
        if 'nexuses/' in path and not path.endswith('NEXUS.md'):
            # 提取可能的 nexus id
            import re
            nexus_match = re.search(r'nexuses/([^/]+)', path)
            nexus_id = nexus_match.group(1) if nexus_match else 'your-nexus-id'
            
            # 如果是写入 .json 或其他配置文件，返回警告并引导正确格式
            if path.endswith('.json') or (path.endswith('.md') and 'NEXUS.md' not in path):
                return (f"【格式提示】检测到你正在向 nexuses/ 目录写入非标准文件。\n\n"
                        f"⚠️ Nexus 只能通过 NEXUS.md 文件定义，系统不会识别 .json 或其他 .md 文件！\n\n"
                        f"📝 正确做法：请创建 nexuses/{nexus_id}/NEXUS.md 文件，格式如下：\n"
                        f"```markdown\n"
                        f"---\n"
                        f"name: Nexus名称\n"
                        f"description: 功能描述\n"
                        f"version: 1.0.0\n"
                        f"skill_dependencies:\n"
                        f"  - 技能ID\n"
                        f"tags:\n"
                        f"  - 标签\n"
                        f"triggers:\n"
                        f"  - 触发词\n"
                        f"objective: 核心目标\n"
                        f"metrics:\n"
                        f"  - 质量指标\n"
                        f"strategy: 执行策略\n"
                        f"---\n\n"
                        f"# Nexus名称 SOP\n\n"
                        f"（详细的标准作业程序）\n"
                        f"```\n\n"
                        f"请使用正确格式重新创建 nexuses/{nexus_id}/NEXUS.md")
        
        # 确保父目录存在
        file_path.parent.mkdir(parents=True, exist_ok=True)
        
        file_path.write_text(content, encoding='utf-8')
        
        # 返回结构化数据，包含完整路径以便前端快速访问
        return json.dumps({
            'action': 'file_created',
            'message': f'已成功写入 {len(content)} 字节',
            'fileName': file_path.name,
            'filePath': str(file_path.resolve()),
            'fileSize': len(content),
        }, ensure_ascii=False)
    
    def _tool_open_in_explorer(self, args: dict) -> str:
        """在文件管理器中打开指定路径并高亮文件"""
        path = args.get('path', '')
        if not path:
            raise ValueError("路径参数不能为空")
        
        file_path = Path(path)
        if not file_path.exists():
            raise FileNotFoundError(f"文件不存在: {path}")
        
        import platform
        import subprocess
        system = platform.system()
        
        try:
            if system == 'Windows':
                # Windows: 使用 explorer /select 高亮文件
                subprocess.run(['explorer', '/select,', str(file_path.resolve())], check=False)
            elif system == 'Darwin':  # macOS
                subprocess.run(['open', '-R', str(file_path.resolve())], check=True)
            else:  # Linux
                # 打开父目录
                subprocess.run(['xdg-open', str(file_path.parent.resolve())], check=True)
            
            return f"已在文件管理器中打开: {file_path.name}"
        except Exception as e:
            raise RuntimeError(f"无法打开文件管理器: {str(e)}")
    
    def _check_nexus_duplication(self, new_content: str) -> str | None:
        """检查新建的 Nexus 是否与现存 Nexus 重复，返回重复的 Nexus ID"""
        # 1. 提取新 Nexus 的 frontmatter
        match = re.match(r'^---\s*\n(.*?)\n---\s*\n', new_content, re.DOTALL)
        if not match:
            return None
        
        new_meta = {}
        if HAS_YAML:
            try:
                new_meta = yaml.safe_load(match.group(1)) or {}
            except Exception:
                pass
        else:
            for line in match.group(1).split('\n'):
                m = re.match(r'^(\w+)\s*:\s*(.+)$', line.strip())
                if m:
                    new_meta[m.group(1)] = m.group(2).strip()
        
        new_name = str(new_meta.get('name', ''))
        new_desc = str(new_meta.get('description', ''))
        if not new_name and not new_desc:
            return None
        
        new_text = f"{new_name} {new_desc}"
        
        # 2. 遍历现有 Nexus 进行对比
        nexuses_dir = self.clawd_path / 'nexuses'
        if not nexuses_dir.exists():
            return None
        
        best_match = None
        highest_score = 0.0
        
        for nexus_md in nexuses_dir.rglob('NEXUS.md'):
            existing_meta = parse_nexus_frontmatter(nexus_md)
            ext_name = str(existing_meta.get('name', ''))
            ext_desc = str(existing_meta.get('description', ''))
            
            ext_text = f"{ext_name} {ext_desc}"
            score = calculate_text_similarity(new_text, ext_text)
            
            if score > highest_score:
                highest_score = score
                best_match = nexus_md.parent.name
        
        # 阈值：超过 55% 的特征重合即判定为重复
        if highest_score >= 0.55:
            return best_match
        
        return None
    
    def _tool_append_file(self, args: dict) -> str:
        """追加内容到文件"""
        path = args.get('path', '')
        content = args.get('content', '')
        
        file_path = self._resolve_path(path)
        file_path.parent.mkdir(parents=True, exist_ok=True)
        
        with open(file_path, 'a', encoding='utf-8') as f:
            f.write(content)
        
        return f"Appended {len(content)} bytes to {file_path.name}"
    
    def _tool_list_dir(self, args: dict) -> str:
        """列出目录内容"""
        path = args.get('path', '.')
        dir_path = self._resolve_path(path)
        
        if not dir_path.exists():
            raise FileNotFoundError(f"Directory not found: {path}")
        if not dir_path.is_dir():
            raise ValueError(f"Not a directory: {path}")
        
        items = []
        for item in sorted(dir_path.iterdir()):
            item_type = 'dir' if item.is_dir() else 'file'
            size = item.stat().st_size if item.is_file() else 0
            items.append({
                'name': item.name,
                'type': item_type,
                'size': size
            })
        
        return json.dumps(items, ensure_ascii=False)
    
    def _tool_run_cmd(self, args: dict) -> str:
        """执行 Shell 命令 (⚠️ 高危操作)"""
        command = args.get('command', '')
        cwd = args.get('cwd', str(self.clawd_path))
        timeout = min(args.get('timeout', 60), 300)  # 最大 5 分钟
        
        if not command:
            raise ValueError("Command cannot be empty")
        
        # 安全检查
        cmd_lower = command.lower()
        for dangerous in DANGEROUS_COMMANDS:
            if dangerous in cmd_lower:
                raise PermissionError(f"Dangerous command blocked: {command}")
        
        try:
            process = subprocess.run(
                command,
                shell=True,
                cwd=cwd,
                capture_output=True,
                text=True,
                timeout=timeout
            )
            
            stdout = process.stdout[:MAX_OUTPUT_SIZE] if process.stdout else ''
            stderr = process.stderr[:MAX_OUTPUT_SIZE] if process.stderr else ''
            
            result_parts = []
            if stdout:
                result_parts.append(f"STDOUT:\n{stdout}")
            if stderr:
                result_parts.append(f"STDERR:\n{stderr}")
            
            rc = process.returncode
            if rc == 0:
                result_parts.append(f"Exit Code: 0 (成功)")
            else:
                # 提供常见 exit code 的可读解释
                code_hints = {
                    1: "通用错误",
                    2: "参数错误或命令误用",
                    3: "URL 格式错误 (curl)",
                    6: "无法解析主机名 (DNS 失败)",
                    7: "无法连接到服务器",
                    28: "操作超时",
                    35: "SSL/TLS 连接错误",
                    56: "网络数据接收失败",
                    60: "SSL 证书验证失败",
                    127: "命令未找到",
                    128: "无效的退出参数",
                }
                hint = code_hints.get(rc, "未知错误")
                result_parts.append(f"Exit Code: {rc} ({hint})")
                # 当没有任何输出时，补充提示帮助 LLM 理解错误
                if not stdout and not stderr:
                    result_parts.append(f"注意: 命令 '{command[:80]}' 执行失败且无输出。建议换用其他工具或方式完成任务。")
            
            return '\n'.join(result_parts)
        
        except subprocess.TimeoutExpired:
            return f"Command timed out after {timeout}s"
    
    def _tool_weather(self, args: dict) -> str:
        """查询天气 (基于 OpenClaw weather skill)"""
        import urllib.request
        import urllib.parse
        
        location = args.get('location', args.get('city', ''))
        if not location:
            raise ValueError("Location/city is required")
        
        # 使用 wttr.in API (无需 API Key)
        encoded_location = urllib.parse.quote(location)
        
        try:
            # 获取详细天气信息
            url = f"https://wttr.in/{encoded_location}?format=j1"
            req = urllib.request.Request(url, headers={'User-Agent': 'curl/7.68.0'})
            
            with urllib.request.urlopen(req, timeout=10) as response:
                data = json.loads(response.read().decode('utf-8'))
            
            current = data.get('current_condition', [{}])[0]
            area = data.get('nearest_area', [{}])[0]
            
            # 格式化输出
            city_name = area.get('areaName', [{}])[0].get('value', location)
            country = area.get('country', [{}])[0].get('value', '')
            
            result = f"""天气查询结果 - {city_name}, {country}

当前温度: {current.get('temp_C', 'N/A')}°C (体感: {current.get('FeelsLikeC', 'N/A')}°C)
天气状况: {current.get('weatherDesc', [{}])[0].get('value', 'N/A')}
湿度: {current.get('humidity', 'N/A')}%
风速: {current.get('windspeedKmph', 'N/A')} km/h ({current.get('winddir16Point', '')})
能见度: {current.get('visibility', 'N/A')} km
紫外线指数: {current.get('uvIndex', 'N/A')}
"""
            return result
            
        except Exception as e:
            # 降级方案：使用简单格式
            try:
                simple_url = f"https://wttr.in/{encoded_location}?format=%l:+%c+%t+(%f)+%h+%w"
                req = urllib.request.Request(simple_url, headers={'User-Agent': 'curl/7.68.0'})
                with urllib.request.urlopen(req, timeout=10) as response:
                    return response.read().decode('utf-8')
            except:
                return f"无法查询 {location} 的天气: {str(e)}"
    
    def _tool_web_search(self, args: dict) -> str:
        """网页搜索 (使用 DuckDuckGo HTML)"""
        import urllib.request
        import urllib.parse
        import re
        
        query = args.get('query', args.get('q', ''))
        if not query:
            raise ValueError("Search query is required")
        
        encoded_query = urllib.parse.quote(query)
        
        try:
            # 使用 DuckDuckGo HTML 版本
            url = f"https://html.duckduckgo.com/html/?q={encoded_query}"
            req = urllib.request.Request(url, headers={
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
            })
            
            with urllib.request.urlopen(req, timeout=15) as response:
                html = response.read().decode('utf-8')
            
            # 提取搜索结果
            results = []
            # 匹配结果链接和标题
            pattern = r'<a[^>]+class="result__a"[^>]*href="([^"]*)"[^>]*>([^<]*)</a>'
            matches = re.findall(pattern, html)
            
            for i, (link, title) in enumerate(matches[:5]):  # 取前5个结果
                # 清理 DuckDuckGo 重定向链接
                if 'uddg=' in link:
                    actual_link = urllib.parse.unquote(link.split('uddg=')[-1].split('&')[0])
                else:
                    actual_link = link
                results.append(f"{i+1}. {title.strip()}\n   {actual_link}")
            
            if results:
                return f"搜索 '{query}' 的结果:\n\n" + "\n\n".join(results)
            else:
                return f"未找到 '{query}' 的相关结果"
                
        except Exception as e:
            return f"搜索失败: {str(e)}"
    
    def _tool_web_fetch(self, args: dict) -> str:
        """获取网页内容 (简化版，提取主要文本)"""
        import urllib.request
        import urllib.parse
        import re
        from html.parser import HTMLParser
        
        url = args.get('url', '')
        if not url:
            raise ValueError("URL is required")
        
        # 确保 URL 有协议
        if not url.startswith(('http://', 'https://')):
            url = 'https://' + url
        
        try:
            req = urllib.request.Request(url, headers={
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36',
                'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8',
            })
            
            with urllib.request.urlopen(req, timeout=15) as response:
                # 检查内容类型
                content_type = response.headers.get('Content-Type', '')
                if 'text/html' not in content_type and 'text/plain' not in content_type:
                    return f"无法读取此类型的内容: {content_type}"
                
                html = response.read().decode('utf-8', errors='ignore')
            
            # 简单的 HTML 文本提取
            # 移除 script 和 style 标签
            html = re.sub(r'<script[^>]*>[\s\S]*?</script>', '', html, flags=re.IGNORECASE)
            html = re.sub(r'<style[^>]*>[\s\S]*?</style>', '', html, flags=re.IGNORECASE)
            html = re.sub(r'<head[^>]*>[\s\S]*?</head>', '', html, flags=re.IGNORECASE)
            
            # 提取 title
            title_match = re.search(r'<title[^>]*>([^<]*)</title>', html, re.IGNORECASE)
            title = title_match.group(1).strip() if title_match else ''
            
            # 移除所有 HTML 标签
            text = re.sub(r'<[^>]+>', ' ', html)
            # 清理多余空白
            text = re.sub(r'\s+', ' ', text).strip()
            # 限制长度
            text = text[:4000]
            
            result = f"URL: {url}\n"
            if title:
                result += f"标题: {title}\n"
            result += f"\n内容摘要:\n{text}"
            
            return result
            
        except urllib.error.HTTPError as e:
            return f"HTTP 错误 {e.code}: {e.reason}"
        except urllib.error.URLError as e:
            return f"无法访问 URL: {e.reason}"
        except Exception as e:
            return f"获取网页失败: {str(e)}"
    
    def _tool_save_memory(self, args: dict) -> str:
        """保存记忆到文件"""
        key = args.get('key', '')
        content = args.get('content', '')
        memory_type = args.get('type', 'general')
        
        if not key or not content:
            raise ValueError("key 和 content 参数必填")
        
        # 记忆存储在 memory 目录下
        memory_dir = self.clawd_path / 'memory'
        memory_dir.mkdir(parents=True, exist_ok=True)
        
        # 按日期组织记忆文件
        today = datetime.now().strftime('%Y-%m-%d')
        memory_file = memory_dir / f'{today}.md'
        
        # 格式化记忆条目
        timestamp = datetime.now().strftime('%H:%M:%S')
        entry = f"\n## [{timestamp}] {key}\n- **类型**: {memory_type}\n- **内容**: {content}\n"
        
        # 追加到记忆文件
        with open(memory_file, 'a', encoding='utf-8') as f:
            f.write(entry)
        
        return f"记忆已保存: {key} (类型: {memory_type})"
    
    def _tool_search_memory(self, args: dict) -> str:
        """检索历史记忆"""
        query = args.get('query', '')
        
        if not query:
            raise ValueError("query 参数必填")
        
        memory_dir = self.clawd_path / 'memory'
        if not memory_dir.exists():
            return "记忆库为空，暂无历史记忆。"
        
        results = []
        query_lower = query.lower()
        
        # 遍历所有记忆文件
        for memory_file in sorted(memory_dir.glob('*.md'), reverse=True)[:7]:  # 最近7天
            try:
                content = memory_file.read_text(encoding='utf-8')
                
                # 按条目分割
                entries = content.split('\n## ')
                for entry in entries:
                    if query_lower in entry.lower():
                        # 提取日期和内容
                        date = memory_file.stem
                        results.append(f"[{date}] {entry.strip()[:200]}")
                        
                        if len(results) >= 5:  # 最多返回5条
                            break
            except Exception:
                continue
            
            if len(results) >= 5:
                break
        
        if results:
            return f"找到 {len(results)} 条相关记忆:\n\n" + "\n\n---\n\n".join(results)
        else:
            return f"未找到与 '{query}' 相关的记忆。"
    
    def _tool_nexus_bind_skill(self, args: dict) -> str:
        """为 Nexus 绑定新技能"""
        nexus_id = args.get('nexusId', '')
        skill_id = args.get('skillId', '')
        if not nexus_id or not skill_id:
            raise ValueError('Missing nexusId or skillId')

        nexus_md = self.clawd_path / 'nexuses' / nexus_id / 'NEXUS.md'
        if not nexus_md.exists():
            raise ValueError(f"Nexus '{nexus_id}' not found")

        # 验证技能存在 (skills/ 目录中有对应目录)
        skill_dir = self.clawd_path / 'skills' / skill_id
        if not skill_dir.exists():
            raise ValueError(f"Skill '{skill_id}' not found in skills/")

        frontmatter = parse_nexus_frontmatter(nexus_md)
        deps = list(frontmatter.get('skill_dependencies', []))

        if skill_id in deps:
            return f"Skill '{skill_id}' already bound to Nexus '{nexus_id}'"

        deps.append(skill_id)
        update_nexus_frontmatter(nexus_md, {'skill_dependencies': deps})
        return f"Skill '{skill_id}' bound to Nexus '{nexus_id}'. Dependencies: {deps}"

    def _tool_nexus_unbind_skill(self, args: dict) -> str:
        """从 Nexus 解绑技能"""
        nexus_id = args.get('nexusId', '')
        skill_id = args.get('skillId', '')
        if not nexus_id or not skill_id:
            raise ValueError('Missing nexusId or skillId')

        nexus_md = self.clawd_path / 'nexuses' / nexus_id / 'NEXUS.md'
        if not nexus_md.exists():
            raise ValueError(f"Nexus '{nexus_id}' not found")

        frontmatter = parse_nexus_frontmatter(nexus_md)
        deps = list(frontmatter.get('skill_dependencies', []))

        if skill_id not in deps:
            return f"Skill '{skill_id}' not bound to Nexus '{nexus_id}'"

        if len(deps) <= 1:
            return f"Cannot remove last skill from Nexus '{nexus_id}'. At least 1 skill required."

        deps.remove(skill_id)
        update_nexus_frontmatter(nexus_md, {'skill_dependencies': deps})
        return f"Skill '{skill_id}' unbound from Nexus '{nexus_id}'. Remaining: {deps}"

    def _tool_generate_skill(self, args: dict) -> str:
        """动态生成 Python SKILL 并保存
        
        当遇到无法完成的任务时，Agent 可以调用此工具生成新的 Python 技能来解决问题。
        生成的技能会保存到 skills/ 目录（或 nexuses/{nexusId}/ 目录）并自动热加载。
        
        参数:
        - name: 技能名称 (kebab-case, 如 "pdf-merger")
        - description: 技能描述
        - pythonCode: Python 实现代码 (必须包含 main() 函数)
        - nexusId: 可选，如果指定则保存到对应 Nexus 目录
        - triggers: 可选，触发关键词列表
        """
        name = args.get('name', '')
        description = args.get('description', '')
        python_code = args.get('pythonCode', '')
        nexus_id = args.get('nexusId', '')
        triggers = args.get('triggers', [])
        
        if not name or not description or not python_code:
            raise ValueError("Missing required parameters: name, description, pythonCode")
        
        # 规范化技能名称 (kebab-case)
        safe_name = re.sub(r'[^\w-]', '-', name.lower()).strip('-')
        safe_name = re.sub(r'-+', '-', safe_name)
        
        if not safe_name:
            raise ValueError("Invalid skill name")
        
        # 验证 Python 代码包含 main() 函数
        if 'def main(' not in python_code and 'async def main(' not in python_code:
            raise ValueError("Python code must contain a main() function")
        
        # 确定保存路径
        if nexus_id:
            # 保存到 Nexus 专属目录
            skill_dir = self.clawd_path / 'nexuses' / nexus_id / 'skills' / safe_name
        else:
            # 保存到全局 skills 目录
            skill_dir = self.clawd_path / 'skills' / safe_name
        
        skill_dir.mkdir(parents=True, exist_ok=True)
        
        # 生成 SKILL.md
        trigger_list = '\n'.join(f'- {t}' for t in triggers) if triggers else f'- {safe_name}'
        skill_md_content = f'''---
name: {safe_name}
description: {description}
version: "1.0.0"
author: auto-generated
triggers:
{trigger_list}
---

# {name}

{description}

## 使用方法

此技能由 DD-OS Agent 自动生成，用于解决特定任务。

### 执行

```bash
python {safe_name}.py
```

### 参数

请参考 Python 代码中的 `main()` 函数签名。

## 实现

参见 `{safe_name}.py`
'''
        
        # 写入文件
        skill_md_path = skill_dir / 'SKILL.md'
        skill_md_path.write_text(skill_md_content, encoding='utf-8')
        
        python_file_path = skill_dir / f'{safe_name}.py'
        python_file_path.write_text(python_code, encoding='utf-8')
        
        # 热加载: 重新注册工具
        try:
            tool_registry.refresh_skills()
            loaded_msg = "并已热加载到工具列表"
        except Exception as e:
            loaded_msg = f"但热加载失败: {e}"
        
        return json.dumps({
            'action': 'skill_created',
            'message': f'技能 "{safe_name}" 已成功创建{loaded_msg}',
            'skillName': safe_name,
            'skillDir': str(skill_dir),
            'files': [str(skill_md_path), str(python_file_path)],
            'nexusId': nexus_id or None,
        }, ensure_ascii=False)

    # ============================================
    # 原有处理器 (保持兼容)
    # ============================================
    
    def handle_index(self):
        html = f"""<!DOCTYPE html>
<html>
<head><title>DD-OS Native Server</title></head>
<body style="font-family: monospace; background: #0f172a; color: #e2e8f0; padding: 30px;">
<h1>DD-OS Native Server v{VERSION}</h1>
<p style="color: #94a3b8;">独立运行的本地 AI 操作系统后端</p>
<p>Clawd Path: <code style="color: #22d3ee;">{self.clawd_path}</code></p>

<h2>📡 API Endpoints</h2>
<div style="background: #1e293b; padding: 15px; border-radius: 8px;">
<h3 style="color: #f59e0b;">数据读取</h3>
<ul>
<li><a href="/status" style="color: #60a5fa;">/status</a> - 服务状态</li>
<li><a href="/files" style="color: #60a5fa;">/files</a> - 文件列表</li>
<li><a href="/file/SOUL.md" style="color: #60a5fa;">/file/SOUL.md</a> - 读取 SOUL</li>
<li><a href="/skills" style="color: #60a5fa;">/skills</a> - 技能列表</li>
<li><a href="/all" style="color: #60a5fa;">/all</a> - 所有数据</li>
</ul>

<h3 style="color: #10b981;">🛠️ 工具执行 (POST)</h3>
<ul>
<li><code>/api/tools/execute</code> - 执行工具</li>
<li>支持: readFile, writeFile, listDir, runCmd, appendFile</li>
</ul>
</div>

<h2>🧪 测试</h2>
<pre style="background: #1e293b; padding: 15px; border-radius: 8px; overflow-x: auto;">
curl -X POST http://localhost:3001/api/tools/execute \\
  -H "Content-Type: application/json" \\
  -d '{{"name": "listDir", "args": {{"path": "."}}}}'
</pre>
</body>
</html>"""
        
        self.send_response(200)
        self.send_header('Content-Type', 'text/html; charset=utf-8')
        self.send_cors_headers()
        self.end_headers()
        self.wfile.write(html.encode('utf-8'))
    
    def handle_status(self):
        files = list_files(self.clawd_path)
        skills_dir = self.clawd_path / 'skills'
        skill_count = len(list(skills_dir.iterdir())) if skills_dir.exists() else 0
        
        self.send_json({
            'status': 'ok',
            'version': VERSION,
            'mode': 'native',
            'clawdPath': str(self.clawd_path),
            'fileCount': len(files),
            'skillCount': skill_count,
            'tools': [t['name'] for t in self.registry.list_all()],
            'toolCount': len(self.registry.list_all()),
            'timestamp': datetime.now().isoformat()
        })
    
    def handle_files(self):
        files = list_files(self.clawd_path)
        self.send_json(files)
    
    def handle_file(self, filename):
        filepath = self.clawd_path / filename
        if not filepath.exists():
            self.send_error_json(f'File not found: {filename}', 404)
            return
        
        if not filepath.is_file():
            self.send_error_json(f'Not a file: {filename}', 400)
            return
        
        try:
            filepath.resolve().relative_to(self.clawd_path.resolve())
        except ValueError:
            self.send_error_json('Access denied', 403)
            return
        
        try:
            content = filepath.read_text(encoding='utf-8')
            self.send_text(content)
        except Exception as e:
            self.send_error_json(f'Read error: {str(e)}', 500)
    
    def handle_skills(self):
        """GET /skills - 递归扫描所有技能 (SKILL.md + manifest.json)，支持用户目录 + 项目目录"""
        skills = []
        seen = set()
        seen_ids = set()  # 防止重复技能 (用户目录优先)
        
        # 获取技能目录列表：用户目录优先，项目目录作为后备
        skills_dirs = []
        user_skills_dir = self.clawd_path / 'skills'
        if user_skills_dir.exists() and user_skills_dir.is_dir():
            skills_dirs.append(('user', user_skills_dir))
        
        project_path = self.project_path or Path(__file__).parent.resolve()
        project_skills_dir = project_path / 'skills'
        if project_skills_dir.exists() and project_skills_dir.is_dir() and project_skills_dir != user_skills_dir:
            skills_dirs.append(('bundled', project_skills_dir))
        
        if not skills_dirs:
            self.send_json([])
            return

        for source, skills_dir in skills_dirs:
            # Phase 1: 扫描有 SKILL.md 的目录
            for skill_md in skills_dir.rglob('SKILL.md'):
                skill_dir = skill_md.parent
                dir_key = str(skill_dir.resolve())
                skill_id = skill_dir.name
                
                if dir_key in seen or skill_id in seen_ids:
                    continue
                seen.add(dir_key)
                seen_ids.add(skill_id)

                frontmatter = parse_skill_frontmatter(skill_md)
                manifest_path = skill_dir / 'manifest.json'

                skill_data = {
                    'id': skill_id,
                    'name': frontmatter.get('name', skill_dir.name),
                    'description': frontmatter.get('description', ''),
                    'location': source,  # 'user' 或 'bundled'
                    'path': str(skill_dir),
                    'status': 'active',
                    'enabled': True,
                    'keywords': frontmatter.get('tags', frontmatter.get('keywords', [])),
                }

                # 无 frontmatter description 时提取正文首段
                if not skill_data['description']:
                    try:
                        content = skill_md.read_text(encoding='utf-8')
                        for line in content.split('\n'):
                            line = line.strip()
                            if line and not line.startswith('#') and not line.startswith('---'):
                                skill_data['description'] = line[:200]
                                break
                    except Exception:
                        pass

                self._enrich_skill_from_manifest(skill_data, manifest_path, frontmatter)
                skills.append(skill_data)

            # Phase 2: 扫描有 manifest.json 但没有 SKILL.md 的目录
            for manifest_path in skills_dir.rglob('manifest.json'):
                skill_dir = manifest_path.parent
                dir_key = str(skill_dir.resolve())
                skill_id = skill_dir.name
                
                if dir_key in seen or skill_id in seen_ids:
                    continue
                seen.add(dir_key)
                seen_ids.add(skill_id)

                try:
                    manifest = json.loads(manifest_path.read_text(encoding='utf-8'))
                except Exception:
                    continue

                skill_data = {
                    'id': skill_id,
                    'name': manifest.get('name', skill_dir.name),
                    'description': manifest.get('description', ''),
                    'location': source,
                    'path': str(skill_dir),
                    'status': 'active',
                    'enabled': True,
                    'keywords': manifest.get('keywords', []),
                }

                self._enrich_skill_from_manifest(skill_data, manifest_path, {})
                skills.append(skill_data)

        self.send_json(skills)

    def _enrich_skill_from_manifest(self, skill_data: dict, manifest_path: Path, frontmatter: dict):
        """从 manifest.json 补充技能元数据 (多工具格式 + 关键词合并)"""
        if manifest_path.exists():
            try:
                manifest = json.loads(manifest_path.read_text(encoding='utf-8'))

                # 支持新格式 tools: [...] 和旧格式 toolName: "..."
                tools_list = manifest.get('tools', [])
                if not tools_list:
                    tools_list = [manifest]

                tool_names = [t.get('toolName') for t in tools_list if t.get('toolName')]
                if tool_names:
                    skill_data['toolNames'] = tool_names
                    skill_data['toolName'] = tool_names[0]
                    skill_data['executable'] = True

                # 合并关键词 (manifest 优先)
                manifest_keywords = list(manifest.get('keywords', []))
                for t in tools_list:
                    manifest_keywords.extend(t.get('keywords', []))
                if manifest_keywords:
                    skill_data['keywords'] = list(set(manifest_keywords))

                skill_data['dangerLevel'] = manifest.get('dangerLevel', 'safe')
                skill_data['version'] = manifest.get('version', '1.0.0')
                if manifest.get('description'):
                    skill_data['description'] = manifest['description']

                # 合并所有 inputs
                all_inputs = {}
                for t in tools_list:
                    all_inputs.update(t.get('inputs', {}))
                if all_inputs:
                    skill_data['inputs'] = all_inputs

            except Exception:
                pass
        else:
            # 纯 SKILL.md - 指令型技能
            skill_data['toolType'] = 'instruction'
            tool_name = skill_name_to_tool_name(skill_data['name'])
            skill_data['toolName'] = tool_name
            skill_data['toolNames'] = [tool_name]
            if frontmatter.get('inputs'):
                skill_data['inputs'] = frontmatter['inputs']

    # ============================================
    # 🌌 Nexus 管理
    # ============================================

    def handle_nexuses(self):
        """GET /nexuses - 扫描 nexuses/ 目录，返回所有 Nexus 列表"""
        nexuses = []
        nexuses_dir = self.clawd_path / 'nexuses'

        if not nexuses_dir.exists():
            nexuses_dir.mkdir(parents=True, exist_ok=True)
            self.send_json([])
            return

        seen = set()

        for nexus_md in nexuses_dir.rglob('NEXUS.md'):
            nexus_dir = nexus_md.parent
            dir_key = str(nexus_dir.resolve())
            if dir_key in seen:
                continue
            seen.add(dir_key)

            frontmatter = parse_nexus_frontmatter(nexus_md)
            if not frontmatter or not frontmatter.get('name'):
                continue

            sop_content = extract_nexus_body(nexus_md)
            exp_dir = nexus_dir / 'experience'
            xp = count_experience_entries(exp_dir) if exp_dir.exists() else 0

            visual_dna = frontmatter.get('visual_dna', {})

            nexus_data = {
                'id': frontmatter.get('name', nexus_dir.name),
                'name': frontmatter.get('name', nexus_dir.name),
                'description': frontmatter.get('description', ''),
                'archetype': frontmatter.get('archetype', 'REACTOR'),
                'version': frontmatter.get('version', '1.0.0'),
                'skillDependencies': frontmatter.get('skill_dependencies', []),
                'tags': frontmatter.get('tags', []),
                'triggers': frontmatter.get('triggers', []),
                'visualDNA': visual_dna,
                'sopContent': sop_content,
                'xp': xp,
                'location': 'local',
                'path': str(nexus_dir),
                'status': 'active',
                # 目标函数驱动字段 (Objective-Driven Execution)
                'objective': frontmatter.get('objective', ''),
                'metrics': frontmatter.get('metrics', []),
                'strategy': frontmatter.get('strategy', ''),
            }
            nexuses.append(nexus_data)

        self.send_json(nexuses)

    def handle_nexuses_health(self):
        """GET /nexuses/health - 检查 nexuses 目录的配置健康状况"""
        nexuses_dir = self.clawd_path / 'nexuses'
        issues = []
        suggestions = []
        stats = {
            'valid_nexuses': 0,
            'orphan_files': 0,
            'missing_nexus_md': 0,
            'invalid_frontmatter': 0,
        }

        if not nexuses_dir.exists():
            self.send_json({
                'healthy': True,
                'issues': [],
                'suggestions': ['nexuses 目录为空，可以开始创建 Nexus'],
                'stats': stats
            })
            return

        # 收集所有有效的 Nexus 目录
        valid_dirs = set()
        for nexus_md in nexuses_dir.rglob('NEXUS.md'):
            nexus_dir = nexus_md.parent
            frontmatter = parse_nexus_frontmatter(nexus_md)
            if frontmatter and frontmatter.get('name'):
                valid_dirs.add(str(nexus_dir.resolve()))
                stats['valid_nexuses'] += 1
            else:
                stats['invalid_frontmatter'] += 1
                issues.append({
                    'type': 'invalid_frontmatter',
                    'path': str(nexus_md),
                    'message': f"NEXUS.md 缺少必要的 'name' 字段",
                })

        # 检查孤立文件（有 .json 但没有 NEXUS.md）
        for item in nexuses_dir.iterdir():
            if item.is_file() and item.suffix == '.json':
                # 检查是否有对应的 NEXUS.md 目录
                stem = item.stem.replace('.json', '')
                potential_dir = nexuses_dir / stem
                if not (potential_dir / 'NEXUS.md').exists():
                    stats['orphan_files'] += 1
                    issues.append({
                        'type': 'orphan_json',
                        'path': str(item),
                        'message': f"发现孤立的 JSON 文件，没有对应的 NEXUS.md",
                        'suggestion': f"创建 {stem}/NEXUS.md 或删除此文件",
                    })
                    suggestions.append(
                        f"文件 '{item.name}' 可能是 AI 生成的配置，需要转换为 NEXUS.md 格式才能被系统识别"
                    )

            # 检查目录但没有 NEXUS.md
            if item.is_dir() and not (item / 'NEXUS.md').exists():
                # 检查目录内是否有其他文件
                files = list(item.iterdir())
                if files:
                    stats['missing_nexus_md'] += 1
                    issues.append({
                        'type': 'missing_nexus_md',
                        'path': str(item),
                        'message': f"目录 '{item.name}' 缺少 NEXUS.md 文件",
                        'files': [f.name for f in files[:5]],
                    })

        healthy = len(issues) == 0
        self.send_json({
            'healthy': healthy,
            'issues': issues,
            'suggestions': suggestions,
            'stats': stats,
            'tip': '运行 /nexuses 查看所有有效的 Nexus' if healthy else '请修复上述问题后重新检查',
        })

    def handle_nexus_detail(self, nexus_name: str):
        """GET /nexuses/{name} - 获取单个 Nexus 完整信息"""
        nexuses_dir = self.clawd_path / 'nexuses'
        nexus_dir = nexuses_dir / nexus_name
        nexus_md = nexus_dir / 'NEXUS.md'

        if not nexus_md.exists():
            self.send_error_json(f"Nexus '{nexus_name}' not found", 404)
            return

        frontmatter = parse_nexus_frontmatter(nexus_md)
        sop_content = extract_nexus_body(nexus_md)
        exp_dir = nexus_dir / 'experience'
        xp = count_experience_entries(exp_dir) if exp_dir.exists() else 0

        # 加载最近经验条目
        recent_experiences = []
        for exp_file in ['successes.md', 'failures.md']:
            exp_path = exp_dir / exp_file
            if not exp_path.exists():
                continue
            try:
                content = exp_path.read_text(encoding='utf-8')
                outcome = 'success' if 'success' in exp_file else 'failure'
                entries = content.split('\n### ')
                for entry in entries[1:]:  # skip header
                    entry = entry.strip()
                    if not entry:
                        continue
                    lines = entry.split('\n')
                    title = lines[0].strip() if lines else ''
                    recent_experiences.append({
                        'title': title,
                        'outcome': outcome,
                        'content': '\n'.join(lines[1:]).strip(),
                    })
            except Exception:
                pass

        # 按时间倒序（标题通常包含日期）
        recent_experiences = recent_experiences[-10:][::-1]

        visual_dna = frontmatter.get('visual_dna', {})

        response = {
            'id': frontmatter.get('name', nexus_name),
            'name': frontmatter.get('name', nexus_name),
            'description': frontmatter.get('description', ''),
            'archetype': frontmatter.get('archetype', 'REACTOR'),
            'version': frontmatter.get('version', '1.0.0'),
            'skillDependencies': frontmatter.get('skill_dependencies', []),
            'tags': frontmatter.get('tags', []),
            'triggers': frontmatter.get('triggers', []),
            'visualDNA': visual_dna,
            'sopContent': sop_content,
            'xp': xp,
            'recentExperiences': recent_experiences,
            'location': 'local',
            'path': str(nexus_dir),
            'status': 'active',
            # 目标函数驱动字段 (Objective-Driven Execution)
            'objective': frontmatter.get('objective', ''),
            'metrics': frontmatter.get('metrics', []),
            'strategy': frontmatter.get('strategy', ''),
        }
        self.send_json(response)

    def handle_nexus_update_skills(self, nexus_name: str, data: dict):
        """POST /nexuses/{name}/skills - 更新 Nexus 技能依赖"""
        action = data.get('action', '')  # 'add' or 'remove'
        skill_id = data.get('skillId', '')

        if action not in ('add', 'remove') or not skill_id:
            self.send_error_json('Invalid: need action (add/remove) and skillId', 400)
            return

        nexus_dir = self.clawd_path / 'nexuses' / nexus_name
        nexus_md = nexus_dir / 'NEXUS.md'
        if not nexus_md.exists():
            self.send_error_json(f"Nexus '{nexus_name}' not found", 404)
            return

        frontmatter = parse_nexus_frontmatter(nexus_md)
        deps = list(frontmatter.get('skill_dependencies', []))

        if action == 'add':
            if skill_id not in deps:
                deps.append(skill_id)
        elif action == 'remove':
            if len(deps) <= 1:
                self.send_error_json('Cannot remove last skill dependency', 400)
                return
            if skill_id in deps:
                deps.remove(skill_id)

        update_nexus_frontmatter(nexus_md, {'skill_dependencies': deps})

        self.send_json({
            'status': 'ok',
            'nexusId': nexus_name,
            'skillDependencies': deps,
        })

    def handle_nexus_update_meta(self, nexus_name: str, data: dict):
        """POST /nexuses/{name}/meta - 更新 Nexus 元数据(名称等)"""
        nexus_dir = self.clawd_path / 'nexuses' / nexus_name
        nexus_md = nexus_dir / 'NEXUS.md'
        
        if not nexus_md.exists():
            self.send_error_json(f"Nexus '{nexus_name}' not found", 404)
            return

        new_name = data.get('name', '').strip()
        if not new_name:
            self.send_error_json('Invalid: name is required', 400)
            return
            
        update_nexus_frontmatter(nexus_md, {'name': new_name})

        self.send_json({
            'status': 'ok',
            'nexusId': nexus_name,
            'name': new_name
        })

    def handle_add_experience(self, nexus_name: str, data: dict):
        """POST /nexuses/{name}/experience - 为 Nexus 添加经验记录"""
        nexuses_dir = self.clawd_path / 'nexuses'
        nexus_dir = nexuses_dir / nexus_name

        if not nexus_dir.exists():
            self.send_error_json(f"Nexus '{nexus_name}' not found", 404)
            return

        task = data.get('task', '')
        tools_used = data.get('tools_used', [])
        outcome = data.get('outcome', 'success')
        key_insight = data.get('key_insight', '')

        if not task:
            self.send_error_json('Missing required field: task', 400)
            return

        # 确保 experience 目录存在
        exp_dir = nexus_dir / 'experience'
        exp_dir.mkdir(parents=True, exist_ok=True)

        # 构建 Markdown 条目
        from datetime import datetime
        timestamp = datetime.now().strftime('%Y-%m-%d')
        tool_seq = ' → '.join(tools_used) if tools_used else 'N/A'

        entry = f"\n### [{timestamp}] {task[:80]}\n"
        entry += f"- **Tools**: {tool_seq}\n"
        if key_insight:
            entry += f"- **Insight**: {key_insight}\n"
        entry += "---\n"

        # 追加到对应文件
        target_file = exp_dir / ('successes.md' if outcome == 'success' else 'failures.md')
        try:
            with target_file.open('a', encoding='utf-8') as f:
                f.write(entry)
            self.send_json({'status': 'ok', 'outcome': outcome})
        except Exception as e:
            self.send_error_json(f'Failed to write experience: {str(e)}', 500)

    def handle_tools_list(self):
        """GET /tools - 列出所有已注册的工具"""
        self.send_json(self.registry.list_all())

    def handle_tools_reload(self, data=None):
        """POST /tools/reload - 热重载插件工具"""
        self.registry.plugin_tools.clear()
        self.registry.instruction_tools.clear()
        self.registry.scan_plugins()
        tools = self.registry.list_all()
        self.send_json({
            'status': 'ok',
            'message': f'Reloaded. {len(tools)} tools registered.',
            'tools': tools,
        })

    # ============================================
    # 🔌 MCP 服务器管理
    # ============================================

    def handle_mcp_servers_list(self):
        """GET /mcp/servers - 列出 MCP 服务器状态"""
        if not self.registry.mcp_manager:
            self.send_json({
                'status': 'ok',
                'enabled': False,
                'servers': {},
                'message': 'MCP support not initialized'
            })
            return

        status = self.registry.mcp_manager.get_server_status()
        mcp_tools = [t for t in self.registry.list_all() if t.get('type') == 'mcp']

        self.send_json({
            'status': 'ok',
            'enabled': True,
            'servers': status,
            'toolCount': len(mcp_tools),
            'tools': mcp_tools,
        })

    def handle_mcp_reload(self, data=None):
        """POST /mcp/reload - 重新加载 MCP 服务器"""
        if not HAS_MCP:
            self.send_json({
                'status': 'error',
                'message': 'MCP support not available'
            }, 400)
            return

        # 清理现有 MCP 工具
        self.registry.mcp_tools.clear()
        if self.registry.mcp_manager:
            self.registry.mcp_manager.shutdown_all()

        # 重新扫描
        self.registry.scan_mcp_servers()

        mcp_tools = [t for t in self.registry.list_all() if t.get('type') == 'mcp']
        self.send_json({
            'status': 'ok',
            'message': f'MCP reloaded. {len(mcp_tools)} tool(s) registered.',
            'tools': mcp_tools,
        })

    def handle_mcp_reconnect(self, server_name: str):
        """POST /mcp/servers/{name}/reconnect - 重连 MCP 服务器"""
        if not self.registry.mcp_manager:
            self.send_json({
                'status': 'error',
                'message': 'MCP support not initialized'
            }, 400)
            return

        success = self.registry.mcp_manager.reconnect_server(server_name)

        if success:
            # 更新工具注册
            self.registry.mcp_tools.clear()
            for tool_info in self.registry.mcp_manager.get_all_tools():
                tool_name = tool_info['name']
                if tool_name not in self.registry.builtin_tools and tool_name not in self.registry.plugin_tools:
                    self.registry.mcp_tools[tool_name] = {
                        'name': tool_name,
                        'server': tool_info.get('server', ''),
                        'description': tool_info.get('description', ''),
                        'inputs': tool_info.get('inputs', {}),
                        'dangerLevel': 'safe',
                        'version': '1.0.0',
                    }

            self.send_json({
                'status': 'ok',
                'message': f'Server {server_name} reconnected',
                'server': server_name,
            })
        else:
            self.send_json({
                'status': 'error',
                'message': f'Failed to reconnect server: {server_name}'
            }, 500)

    # ============================================
    # 🔍 Registry 在线搜索 (TF-IDF, 无 LLM)
    # ============================================

    def handle_registry_skills_search(self, query: dict):
        """GET /api/registry/skills?q={query} - 搜索可安装的技能"""
        q = query.get('q', [''])[0].strip().lower()
        
        # 读取 registry 文件
        registry_path = self.clawd_path / 'registry' / 'skills.json'
        if not registry_path.exists():
            self.send_json({'status': 'ok', 'results': [], 'message': 'Registry not found'})
            return
        
        try:
            registry = json.loads(registry_path.read_text(encoding='utf-8'))
            skills = registry.get('skills', [])
        except Exception as e:
            self.send_json({'status': 'error', 'message': f'Failed to read registry: {e}'}, 500)
            return
        
        # 如果没有查询词，返回所有
        if not q:
            self.send_json({
                'status': 'ok',
                'results': skills[:20],
                'total': len(skills)
            })
            return
        
        # TF-IDF 风格的关键词匹配
        tokens = self._tokenize(q)
        scored_results = []
        
        for skill in skills:
            score = self._compute_skill_score(skill, tokens)
            if score > 0:
                scored_results.append({**skill, 'score': score})
        
        # 按分数排序
        scored_results.sort(key=lambda x: x['score'], reverse=True)
        
        self.send_json({
            'status': 'ok',
            'results': scored_results[:10],
            'total': len(scored_results),
            'query': q
        })

    def handle_registry_mcp_search(self, query: dict):
        """GET /api/registry/mcp?q={query} - 搜索可安装的 MCP 服务器"""
        q = query.get('q', [''])[0].strip().lower()
        
        # 读取 registry 文件
        registry_path = self.clawd_path / 'registry' / 'mcp-servers.json'
        if not registry_path.exists():
            self.send_json({'status': 'ok', 'results': [], 'message': 'Registry not found'})
            return
        
        try:
            registry = json.loads(registry_path.read_text(encoding='utf-8'))
            servers = registry.get('servers', [])
        except Exception as e:
            self.send_json({'status': 'error', 'message': f'Failed to read registry: {e}'}, 500)
            return
        
        # 如果没有查询词，返回所有
        if not q:
            self.send_json({
                'status': 'ok',
                'results': servers[:20],
                'total': len(servers)
            })
            return
        
        # TF-IDF 风格的关键词匹配
        tokens = self._tokenize(q)
        scored_results = []
        
        for server in servers:
            score = self._compute_mcp_score(server, tokens)
            if score > 0:
                scored_results.append({**server, 'score': score})
        
        # 按分数排序
        scored_results.sort(key=lambda x: x['score'], reverse=True)
        
        self.send_json({
            'status': 'ok',
            'results': scored_results[:10],
            'total': len(scored_results),
            'query': q
        })

    def _tokenize(self, text: str) -> list:
        """分词：按空格和标点拆分"""
        import re
        tokens = re.split(r'[\s,，.。!！?？、;；:：\-—]+', text.lower())
        return [t for t in tokens if t and len(t) >= 1]

    def _compute_skill_score(self, skill: dict, tokens: list) -> float:
        """计算技能的匹配分数 (TF-IDF 简化版)"""
        score = 0.0
        name = skill.get('name', '').lower()
        desc = skill.get('description', '').lower()
        keywords = [k.lower() for k in skill.get('keywords', [])]
        full_text = f"{name} {desc} {' '.join(keywords)}"
        
        for token in tokens:
            # 词频 (TF)
            tf = full_text.count(token)
            # 长词权重更高 (简化 IDF)
            idf = 1.5 if len(token) > 3 else 1.0
            score += tf * idf
            
            # 精确匹配加权
            if token in name:
                score += 10
            if token in keywords:
                score += 5
        
        return min(score, 100)

    def _compute_mcp_score(self, server: dict, tokens: list) -> float:
        """计算 MCP 服务器的匹配分数"""
        score = 0.0
        name = server.get('name', '').lower()
        desc = server.get('description', '').lower()
        keywords = [k.lower() for k in server.get('keywords', [])]
        full_text = f"{name} {desc} {' '.join(keywords)}"
        
        for token in tokens:
            tf = full_text.count(token)
            idf = 1.5 if len(token) > 3 else 1.0
            score += tf * idf
            
            if token in name:
                score += 10
            if token in keywords:
                score += 5
        
        return min(score, 100)

    def handle_mcp_install(self, data: dict):
        """POST /mcp/install - 安装 MCP 服务器配置"""
        server_id = data.get('id', '')
        server_name = data.get('name', server_id)
        command = data.get('command', '')
        args = data.get('args', [])
        env = data.get('env', {})
        
        if not server_name or not command:
            self.send_error_json('Missing required fields: name, command', 400)
            return
        
        # 安全检查
        if '..' in server_name or '/' in server_name or '\\' in server_name:
            self.send_error_json('Invalid server name', 400)
            return
        
        # 读取现有配置
        config_path = self.clawd_path / 'mcp-servers.json'
        try:
            if config_path.exists():
                config = json.loads(config_path.read_text(encoding='utf-8'))
            else:
                config = {'servers': {}}
        except Exception as e:
            self.send_error_json(f'Failed to read config: {e}', 500)
            return
        
        # 检查是否已存在
        if server_name in config.get('servers', {}):
            self.send_error_json(f'Server already exists: {server_name}', 409)
            return
        
        # 添加新服务器配置
        config['servers'][server_name] = {
            'command': command,
            'args': args,
            'env': env,
            'enabled': False  # 默认禁用，需要用户手动启用
        }
        
        # 写回配置文件
        try:
            config_path.write_text(json.dumps(config, indent=2, ensure_ascii=False), encoding='utf-8')
        except Exception as e:
            self.send_error_json(f'Failed to write config: {e}', 500)
            return
        
        self.send_json({
            'status': 'ok',
            'serverName': server_name,
            'message': f'MCP server "{server_name}" added (disabled by default). Enable it in mcp-servers.json to use.',
            'configPath': str(config_path)
        })

    # ============================================
    # 📦 远程技能安装/卸载
    # ============================================

    def handle_skill_install(self, data):
        """POST /skills/install - 从 Git URL 安装技能"""
        source = data.get('source', '')
        name = data.get('name', '')

        if not source:
            self.send_error_json('Missing source parameter', 400)
            return

        if not source.startswith(('http://', 'https://', 'git@')):
            self.send_error_json('Unsupported source format. Use a Git URL (https://... or git@...)', 400)
            return

        try:
            # 从 URL 提取仓库名
            match = re.search(r'/([^/]+?)(?:\.git)?$', source)
            repo_name = name or (match.group(1) if match else 'downloaded-skill')
            # 安全化名称
            repo_name = re.sub(r'[^\w\-.]', '_', repo_name)

            target = self.clawd_path / 'skills' / repo_name
            if target.exists():
                self.send_error_json(f'Skill already exists: {repo_name}. Use /skills/uninstall first.', 409)
                return

            # 确保 skills/ 目录存在
            (self.clawd_path / 'skills').mkdir(parents=True, exist_ok=True)

            # Git clone (shallow, 限制深度)
            process = subprocess.run(
                ['git', 'clone', '--depth', '1', source, str(target)],
                capture_output=True,
                text=True,
                timeout=120,
            )

            if process.returncode != 0:
                # 清理失败的 clone
                if target.exists():
                    shutil.rmtree(target, ignore_errors=True)
                stderr = process.stderr[:500] if process.stderr else 'Unknown error'
                raise RuntimeError(f"git clone failed: {stderr}")

            # 验证: 必须有 SKILL.md 或 manifest.json
            has_skill_md = (target / 'SKILL.md').exists() or any(target.rglob('SKILL.md'))
            has_manifest = (target / 'manifest.json').exists()

            if not has_skill_md and not has_manifest:
                shutil.rmtree(target, ignore_errors=True)
                self.send_error_json('Invalid skill: no SKILL.md or manifest.json found', 400)
                return

            # 重新扫描注册
            self.registry.plugin_tools.clear()
            self.registry.instruction_tools.clear()
            self.registry.scan_plugins()

            self.send_json({
                'status': 'ok',
                'name': repo_name,
                'path': str(target),
                'message': f'Skill installed: {repo_name}',
                'toolCount': len(self.registry.list_all()),
            })

        except subprocess.TimeoutExpired:
            if target.exists():
                shutil.rmtree(target, ignore_errors=True)
            self.send_error_json('Git clone timed out (120s limit)', 500)
        except RuntimeError as e:
            self.send_error_json(str(e), 500)
        except Exception as e:
            self.send_error_json(f'Installation failed: {str(e)}', 500)

    def handle_skill_uninstall(self, data):
        """POST /skills/uninstall - 卸载技能"""
        name = data.get('name', '')

        if not name:
            self.send_error_json('Missing name parameter', 400)
            return

        # 安全检查: 不允许路径遍历
        if '..' in name or '/' in name or '\\' in name:
            self.send_error_json('Invalid skill name', 400)
            return

        target = self.clawd_path / 'skills' / name

        if not target.exists():
            self.send_error_json(f'Skill not found: {name}', 404)
            return

        if not target.is_dir():
            self.send_error_json(f'Not a directory: {name}', 400)
            return

        try:
            shutil.rmtree(target)

            # 重新扫描
            self.registry.plugin_tools.clear()
            self.registry.instruction_tools.clear()
            self.registry.scan_plugins()

            self.send_json({
                'status': 'ok',
                'name': name,
                'message': f'Skill uninstalled: {name}',
                'toolCount': len(self.registry.list_all()),
            })

        except Exception as e:
            self.send_error_json(f'Uninstall failed: {str(e)}', 500)

    def handle_trace_save(self, data):
        """POST /api/traces/save - 保存执行追踪 (P2: 执行流记忆)"""
        if not data:
            self.send_error_json('Missing trace data', 400)
            return

        traces_dir = self.clawd_path / 'memory' / 'exec_traces'
        traces_dir.mkdir(parents=True, exist_ok=True)

        # 按月分片存储
        month = datetime.now().strftime('%Y-%m')
        trace_file = traces_dir / f'{month}.jsonl'

        # 敏感数据脱敏
        trace_json = json.dumps(data, ensure_ascii=False)
        import re
        trace_json = re.sub(
            r'(password|token|secret|api_key|apikey|auth)["\s:]*["\']([^"\']{3,})["\']',
            r'\1": "***"',
            trace_json,
            flags=re.IGNORECASE
        )

        try:
            with open(trace_file, 'a', encoding='utf-8') as f:
                f.write(trace_json + '\n')

            self.send_json({
                'status': 'ok',
                'message': f'Trace saved to {month}.jsonl',
            })
        except Exception as e:
            self.send_error_json(f'Failed to save trace: {e}', 500)

    def handle_trace_search(self, query_params):
        """GET /api/traces/search?query=xxx&limit=5 - 检索执行追踪 (P2)"""
        query = query_params.get('query', [''])[0]
        limit = min(int(query_params.get('limit', ['5'])[0]), 20)

        if not query:
            self.send_json([])
            return

        traces_dir = self.clawd_path / 'memory' / 'exec_traces'
        if not traces_dir.exists():
            self.send_json([])
            return

        query_lower = query.lower()
        query_words = [w for w in query_lower.split() if len(w) > 1]
        results = []

        # 从最近的月份文件开始搜索
        for trace_file in sorted(traces_dir.glob('*.jsonl'), reverse=True)[:6]:
            try:
                for line in reversed(trace_file.read_text(encoding='utf-8').strip().split('\n')):
                    if not line.strip():
                        continue
                    try:
                        trace = json.loads(line)
                        task = trace.get('task', '').lower()
                        tags = [t.lower() for t in trace.get('tags', [])]
                        # 关键词匹配: task 描述或 tags
                        matched = any(w in task for w in query_words) or \
                                  any(w in ' '.join(tags) for w in query_words)
                        if matched:
                            results.append(trace)
                            if len(results) >= limit:
                                break
                    except json.JSONDecodeError:
                        continue
            except Exception:
                continue
            if len(results) >= limit:
                break

        self.send_json(results)
    
    def handle_trace_recent(self, query_params):
        """GET /api/traces/recent?days=3&limit=100 - 获取最近N天的执行日志 (供 Observer 分析)"""
        days = min(int(query_params.get('days', ['3'])[0]), 30)
        limit = min(int(query_params.get('limit', ['100'])[0]), 500)
        
        traces_dir = self.clawd_path / 'memory' / 'exec_traces'
        if not traces_dir.exists():
            self.send_json({'traces': [], 'stats': {}})
            return
        
        cutoff_time = datetime.now() - timedelta(days=days)
        cutoff_ts = cutoff_time.timestamp() * 1000  # 毫秒时间戳
        
        traces = []
        tool_freq = {}  # 工具使用频率
        nexus_freq = {}  # Nexus 使用频率
        total_turns = 0
        total_errors = 0
        
        # 从最近的月份文件开始读取
        for trace_file in sorted(traces_dir.glob('*.jsonl'), reverse=True)[:3]:
            try:
                for line in reversed(trace_file.read_text(encoding='utf-8').strip().split('\n')):
                    if not line.strip():
                        continue
                    try:
                        trace = json.loads(line)
                        ts = trace.get('timestamp', 0)
                        if ts < cutoff_ts:
                            continue  # 超出时间范围
                        
                        traces.append(trace)
                        
                        # 统计工具频率
                        for tool in trace.get('tools', []):
                            tool_name = tool.get('name', 'unknown')
                            tool_freq[tool_name] = tool_freq.get(tool_name, 0) + 1
                        
                        # 统计 Nexus 频率
                        nexus_id = trace.get('activeNexusId')
                        if nexus_id:
                            nexus_freq[nexus_id] = nexus_freq.get(nexus_id, 0) + 1
                        
                        # 统计轮次和错误
                        total_turns += trace.get('turnCount', 0)
                        total_errors += trace.get('errorCount', 0)
                        
                        if len(traces) >= limit:
                            break
                    except json.JSONDecodeError:
                        continue
            except Exception:
                continue
            if len(traces) >= limit:
                break
        
        # 按时间倒序排列
        traces.sort(key=lambda t: t.get('timestamp', 0), reverse=True)
        
        self.send_json({
            'traces': traces,
            'stats': {
                'totalExecutions': len(traces),
                'toolFrequency': tool_freq,
                'nexusFrequency': nexus_freq,
                'avgTurnsPerExecution': total_turns / len(traces) if traces else 0,
                'totalErrors': total_errors,
                'timeRangeDays': days,
            }
        })
    
    def handle_memories(self):
        memories = []
        
        memory_md = self.clawd_path / 'MEMORY.md'
        if memory_md.exists():
            try:
                content = memory_md.read_text(encoding='utf-8')
                memories.extend(parse_memory_md(content))
            except:
                pass
        
        memory_dir = self.clawd_path / 'memory'
        if memory_dir.exists() and memory_dir.is_dir():
            for item in memory_dir.iterdir():
                if item.is_file() and item.suffix == '.md':
                    try:
                        content = item.read_text(encoding='utf-8')
                        memories.append({
                            'id': f'file-{item.stem}',
                            'title': item.stem.replace('-', ' ').replace('_', ' ').title(),
                            'content': content[:500],
                            'type': 'long-term',
                            'timestamp': item.stat().st_mtime,
                            'tags': [],
                        })
                    except:
                        pass
        
        self.send_json(memories)
    
    def handle_all(self):
        data = {
            'soul': None,
            'identity': None,
            'skills': [],
            'memories': [],
            'files': list_files(self.clawd_path),
        }
        
        soul_path = self.clawd_path / 'SOUL.md'
        if soul_path.exists():
            try:
                data['soul'] = soul_path.read_text(encoding='utf-8')
            except:
                pass
        
        identity_path = self.clawd_path / 'IDENTITY.md'
        if identity_path.exists():
            try:
                data['identity'] = identity_path.read_text(encoding='utf-8')
            except:
                pass
        
        skills_dir = self.clawd_path / 'skills'
        if skills_dir.exists():
            for item in skills_dir.iterdir():
                if item.is_dir():
                    data['skills'].append({
                        'name': item.name,
                        'location': 'local',
                        'status': 'active',
                        'enabled': True,
                    })
        
        memory_md = self.clawd_path / 'MEMORY.md'
        if memory_md.exists():
            try:
                content = memory_md.read_text(encoding='utf-8')
                data['memories'] = parse_memory_md(content)
            except:
                pass
        
        self.send_json(data)
    
    def handle_task_execute(self, data):
        """兼容旧的任务执行接口"""
        prompt = data.get('prompt', '').strip()
        if not prompt:
            self.send_error_json('Missing prompt', 400)
            return
        
        task_id = str(uuid.uuid4())[:8]
        
        thread = threading.Thread(
            target=run_task_in_background,
            args=(task_id, prompt, self.clawd_path),
            daemon=True,
        )
        thread.start()
        
        self.send_json({
            'taskId': task_id,
            'status': 'running',
        })
    
    # ============================================
    # 🤖 子代理 API 处理器 (Quest 模式支持)
    # ============================================
    
    def handle_subagent_spawn(self, data):
        """启动子代理"""
        if not self.subagent_manager:
            self.send_error_json('SubagentManager not initialized', 500)
            return
        
        agent_type = data.get('type', 'explore')
        task = data.get('task', '')
        tools = data.get('tools', [])
        context = data.get('context', '')
        
        if not task:
            self.send_error_json('Missing task', 400)
            return
        
        try:
            agent_id = self.subagent_manager.spawn(agent_type, task, tools, context)
            self.send_json({
                'status': 'success',
                'agentId': agent_id,
                'message': f'Spawned {agent_type} agent'
            })
        except Exception as e:
            self.send_error_json(f'Failed to spawn agent: {e}', 500)
    
    def handle_subagent_status(self, agent_id):
        """获取子代理状态"""
        if not self.subagent_manager:
            self.send_error_json('SubagentManager not initialized', 500)
            return
        
        status = self.subagent_manager.get_status(agent_id)
        if status:
            self.send_json({'status': 'success', 'agent': status})
        else:
            self.send_error_json(f'Agent not found: {agent_id}', 404)
    
    def handle_subagent_collect(self, data):
        """收集多个子代理的结果"""
        if not self.subagent_manager:
            self.send_error_json('SubagentManager not initialized', 500)
            return
        
        agent_ids = data.get('agentIds', [])
        timeout = data.get('timeout', 60.0)
        
        if not agent_ids:
            # 返回所有代理状态
            all_status = self.subagent_manager.get_all_status()
            self.send_json({'status': 'success', 'agents': all_status})
            return
        
        try:
            results = self.subagent_manager.collect_results(agent_ids, timeout)
            self.send_json({'status': 'success', 'results': results})
        except Exception as e:
            self.send_error_json(f'Failed to collect results: {e}', 500)
    
    def handle_task_status(self, task_id, offset=0):
        with self.tasks_lock:
            task = self.tasks.get(task_id)
        
        if not task:
            self.send_error_json(f'Task not found: {task_id}', 404)
            return
        
        log_path = task.get('logPath')
        content = ''
        new_offset = offset
        has_more = False
        file_size = task.get('fileSize', 0)
        
        if log_path:
            content, new_offset, has_more = read_log_chunk(log_path, offset)
            try:
                file_size = Path(log_path).stat().st_size
            except:
                pass
        
        self.send_json({
            'taskId': task_id,
            'status': task['status'],
            'content': content,
            'offset': new_offset,
            'hasMore': has_more,
            'fileSize': file_size,
        })


# ============================================
# 辅助函数
# ============================================

def list_files(clawd_path):
    files = []
    try:
        for item in clawd_path.iterdir():
            if item.is_file():
                files.append(item.name)
    except:
        pass
    return sorted(files)


def parse_memory_md(content):
    memories = []
    sections = content.split('## ')
    
    for i, section in enumerate(sections[1:], 1):
        lines = section.strip().split('\n')
        if not lines:
            continue
        
        title = lines[0].strip()
        body = '\n'.join(lines[1:]).strip()
        
        if title:
            memories.append({
                'id': f'memory-{i}',
                'title': title,
                'content': body[:500] if body else title,
                'type': 'long-term',
                'timestamp': None,
                'tags': [],
            })
    
    return memories


def read_log_chunk(log_path, offset=0, max_bytes=51200):
    path = Path(log_path)
    if not path.exists():
        return ('', offset, False)
    
    try:
        file_size = path.stat().st_size
    except:
        return ('', offset, False)
    
    if offset >= file_size:
        return ('', offset, False)
    
    try:
        with open(path, 'rb') as f:
            f.seek(offset)
            raw = f.read(max_bytes)
        
        content = raw.decode('utf-8', errors='replace')
        new_offset = offset + len(raw)
        has_more = new_offset < file_size
        return (content, new_offset, has_more)
    except Exception as e:
        return (f'[日志读取错误: {e}]', offset, False)


def run_task_in_background(task_id, prompt, clawd_path):
    logs_dir = clawd_path / 'logs'
    logs_dir.mkdir(exist_ok=True)
    log_file = logs_dir / f"{task_id}.log"
    
    with ClawdDataHandler.tasks_lock:
        ClawdDataHandler.tasks[task_id] = {
            'taskId': task_id,
            'status': 'running',
            'logPath': str(log_file),
            'fileSize': 0,
        }
    
    try:
        with open(log_file, 'w', encoding='utf-8') as f:
            f.write(f"Task: {prompt}\n")
            f.write(f"Started: {datetime.now().isoformat()}\n")
            f.write("-" * 50 + "\n\n")
        
        # 尝试运行 clawdbot，如果不存在则模拟
        try:
            with open(log_file, 'ab') as f:
                process = subprocess.Popen(
                    ['clawdbot', 'agent', '--agent', 'main', '--message', prompt],
                    cwd=str(clawd_path),
                    stdout=f,
                    stderr=subprocess.STDOUT,
                )
                
                start_time = time.time()
                timeout = 300
                
                while process.poll() is None:
                    time.sleep(0.5)
                    try:
                        with ClawdDataHandler.tasks_lock:
                            ClawdDataHandler.tasks[task_id]['fileSize'] = log_file.stat().st_size
                    except:
                        pass
                    
                    if time.time() - start_time > timeout:
                        process.kill()
                        process.wait()
                        with ClawdDataHandler.tasks_lock:
                            ClawdDataHandler.tasks[task_id]['status'] = 'error'
                            ClawdDataHandler.tasks[task_id]['fileSize'] = log_file.stat().st_size
                        with open(log_file, 'a', encoding='utf-8') as ef:
                            ef.write(f'\n\n[错误] 任务执行超时 ({timeout}s)\n')
                        return
                
                process.wait()
            
            with ClawdDataHandler.tasks_lock:
                ClawdDataHandler.tasks[task_id]['status'] = 'done' if process.returncode == 0 else 'error'
                ClawdDataHandler.tasks[task_id]['fileSize'] = log_file.stat().st_size
        
        except FileNotFoundError:
            # clawdbot 不存在，使用 Native 模式提示
            with open(log_file, 'a', encoding='utf-8') as f:
                f.write("\n[DD-OS Native] clawdbot 未安装。\n")
                f.write("在 Native 模式下，请使用 /api/tools/execute 接口直接执行工具。\n")
                f.write("\n任务已记录，等待 AI 引擎处理。\n")
            
            with ClawdDataHandler.tasks_lock:
                ClawdDataHandler.tasks[task_id]['status'] = 'done'
                ClawdDataHandler.tasks[task_id]['fileSize'] = log_file.stat().st_size
    
    except Exception as e:
        with open(log_file, 'a', encoding='utf-8') as ef:
            ef.write(f'\n\n[错误] {str(e)}\n')
        with ClawdDataHandler.tasks_lock:
            ClawdDataHandler.tasks[task_id]['status'] = 'error'
            ClawdDataHandler.tasks[task_id]['fileSize'] = log_file.stat().st_size


def cleanup_old_logs(clawd_path, max_age_hours=24):
    logs_dir = clawd_path / 'logs'
    if not logs_dir.exists():
        return
    
    now = time.time()
    count = 0
    for f in logs_dir.glob('*.log'):
        try:
            age = now - f.stat().st_mtime
            if age > max_age_hours * 3600:
                f.unlink()
                count += 1
        except:
            pass
    
    if count > 0:
        print(f"[Cleanup] Removed {count} old log files")


def cleanup_old_traces(clawd_path, max_months=6):
    """清理过期的执行追踪文件 (P2: 保留最近N个月)"""
    traces_dir = clawd_path / 'memory' / 'exec_traces'
    if not traces_dir.exists():
        return

    files = sorted(traces_dir.glob('*.jsonl'))
    if len(files) <= max_months:
        return

    old_files = files[:-max_months]
    for f in old_files:
        try:
            f.unlink()
            print(f"[Cleanup] Removed old trace: {f.name}")
        except:
            pass


def cleanup_temp_uploads(clawd_path, max_age_hours=1):
    """清理超过指定时间的临时上传文件"""
    upload_dir = clawd_path / 'temp' / 'uploads'
    if not upload_dir.exists():
        return
    
    now = time.time()
    count = 0
    for f in upload_dir.iterdir():
        try:
            if f.is_file() and (now - f.stat().st_mtime) > max_age_hours * 3600:
                f.unlink()
                count += 1
        except:
            pass
    
    if count > 0:
        print(f"[Cleanup] Removed {count} old temp upload files")


def main():
    parser = argparse.ArgumentParser(description='DD-OS Native Server')
    parser.add_argument('--port', type=int, default=3001, help='Server port (default: 3001)')
    # 支持环境变量覆盖默认路径
    default_path = os.getenv('DDOS_DATA_PATH', '~/.ddos')
    parser.add_argument('--path', type=str, default=default_path, help='Data directory path (default: ~/.ddos)')
    parser.add_argument('--host', type=str, default='0.0.0.0', help='Server host (default: 0.0.0.0)')
    args = parser.parse_args()
    
    clawd_path = Path(args.path).expanduser().resolve()
    
    if not clawd_path.exists():
        print(f"Creating data directory: {clawd_path}")
        clawd_path.mkdir(parents=True, exist_ok=True)
        
        # 创建默认 SOUL.md
        soul_file = clawd_path / 'SOUL.md'
        soul_file.write_text("""# DD-OS Native Soul

You are DD-OS, a local AI operating system running directly on the user's computer.

## Core Principles
- Be helpful and efficient
- Protect user data and privacy
- Execute tasks safely
- Learn from interactions

## Available Tools
- readFile: Read file contents
- writeFile: Write file contents
- listDir: List directory contents
- runCmd: Execute shell commands

## Safety Rules
- Never delete system files
- Ask before destructive operations
- Keep execution logs
""", encoding='utf-8')
        print(f"Created default SOUL.md")
    
    logs_dir = clawd_path / 'logs'
    logs_dir.mkdir(exist_ok=True)
    
    memory_dir = clawd_path / 'memory'
    memory_dir.mkdir(exist_ok=True)
    
    skills_dir = clawd_path / 'skills'
    skills_dir.mkdir(exist_ok=True)
    
    cleanup_old_logs(clawd_path)
    
    # 🔌 初始化工具注册表
    registry = ToolRegistry(clawd_path)
    # 注册内置工具
    builtin_names = [
        'readFile', 'writeFile', 'appendFile', 'listDir', 'runCmd',
        'weather', 'webSearch', 'webFetch', 'saveMemory', 'searchMemory',
        'nexusBindSkill', 'nexusUnbindSkill', 'openInExplorer', 'parseFile',
        'generateSkill',
    ]
    for name in builtin_names:
        registry.register_builtin(name, name)  # handler resolved at dispatch time
    # 扫描插件工具
    registry.scan_plugins()
    # 扫描 MCP 服务器
    registry.scan_mcp_servers()

    # 清理过期执行追踪 (P2: 保留最近6个月)
    cleanup_old_traces(clawd_path)
    cleanup_temp_uploads(clawd_path)

    # 项目目录 (脚本所在目录)
    project_path = Path(__file__).parent.resolve()
    
    ClawdDataHandler.clawd_path = clawd_path
    ClawdDataHandler.project_path = project_path
    ClawdDataHandler.registry = registry
    ClawdDataHandler.subagent_manager = SubagentManager(registry)
    
    server = ThreadingHTTPServer((args.host, args.port), ClawdDataHandler)
    
    tool_names = [t['name'] for t in registry.list_all()]
    plugin_count = len(registry.plugin_tools)
    mcp_count = len(registry.mcp_tools)
    print(f"""
+==================================================================+
|              DD-OS Native Server v{VERSION}                         |
+==================================================================+
|  Mode:    NATIVE (standalone, no OpenClaw needed)                |
|  Server:  http://{args.host}:{args.port}                                    |
|  Data:    {str(clawd_path)[:50]:<50} |
+------------------------------------------------------------------+
|  Tools:   {len(tool_names)} registered ({len(builtin_names)} builtin + {plugin_count} plugins + {mcp_count} mcp)    |
|  API:     /api/tools/execute (POST)  |  /tools (GET)            |
+==================================================================+
    """)
    
    print(f"Press Ctrl+C to stop\n")
    
    try:
        server.serve_forever()
    except KeyboardInterrupt:
        print("\nShutting down...")
        server.shutdown()


if __name__ == '__main__':
    main()
