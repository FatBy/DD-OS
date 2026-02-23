#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
创建余杭高质量发展PPT内容（专业设计增强版）
基于之前分析提取的余杭高质量发展信息
应用专业设计样式：配色方案、字体、布局、视觉元素
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE
import datetime
import os

def create_professional_ppt():
    """创建专业设计的余杭高质量发展PPT"""
    
    # 创建演示文稿
    prs = Presentation()
    
    # 设置幻灯片宽度和高度（16:9标准宽屏）
    prs.slide_width = Inches(13.33)  # 33.87 cm
    prs.slide_height = Inches(7.5)   # 19.05 cm
    
    # ========== 定义专业配色方案 ==========
    COLOR_SCHEME = {
        'primary_dark': RGBColor(0, 51, 102),      # 深蓝色 - 主色调
        'primary_medium': RGBColor(0, 102, 204),   # 中蓝色 - 次级标题
        'primary_light': RGBColor(102, 178, 255),  # 浅蓝色 - 强调色
        'accent_green': RGBColor(0, 153, 76),      # 绿色 - 成功/积极
        'accent_orange': RGBColor(255, 153, 0),    # 橙色 - 警告/突出
        'accent_red': RGBColor(204, 0, 0),         # 红色 - 重要/紧急
        'dark_text': RGBColor(51, 51, 51),         # 深灰色 - 正文
        'light_text': RGBColor(102, 102, 102),     # 浅灰色 - 辅助文字
        'background': RGBColor(245, 245, 245),     # 浅灰色背景
        'white': RGBColor(255, 255, 255),          # 白色
    }
    
    # 定义字体方案
    FONT_SCHEME = {
        'title_main': ('微软雅黑', 44, True),      # 主标题
        'title_section': ('微软雅黑', 36, True),   # 章节标题
        'title_subsection': ('微软雅黑', 28, True), # 子章节标题
        'content_heading': ('微软雅黑', 20, True), # 内容标题
        'content_main': ('微软雅黑', 18, False),   # 正文内容
        'content_small': ('微软雅黑', 16, False),  # 小字/脚注
        'quote': ('楷体', 20, True),               # 引用文字
    }
    
    # ========== 1. 封面页（专业设计） ==========
    slide_layout = prs.slide_layouts[0]  # 标题幻灯片
    slide = prs.slides.add_slide(slide_layout)
    
    # 设置背景色
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_SCHEME['primary_dark']
    
    # 标题
    title = slide.shapes.title
    title.text = "余杭高质量发展路径与实践"
    
    # 设置标题样式
    title_frame = title.text_frame
    title_frame.paragraphs[0].font.name = FONT_SCHEME['title_main'][0]
    title_frame.paragraphs[0].font.size = Pt(FONT_SCHEME['title_main'][1])
    title_frame.paragraphs[0].font.bold = FONT_SCHEME['title_main'][2]
    title_frame.paragraphs[0].font.color.rgb = COLOR_SCHEME['white']
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 副标题
    subtitle = slide.placeholders[1]
    subtitle.text = f"数字经济时代下的区域发展新范式\n{datetime.datetime.now().strftime('%Y年%m月%d日')}"
    
    # 设置副标题样式
    subtitle_frame = subtitle.text_frame
    subtitle_frame.paragraphs[0].font.name = '微软雅黑'
    subtitle_frame.paragraphs[0].font.size = Pt(24)
    subtitle_frame.paragraphs[0].font.color.rgb = COLOR_SCHEME['primary_light']
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 添加装饰线条
    left = Inches(2)
    top = Inches(5.5)
    width = Inches(9.33)
    height = Pt(2)
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_SCHEME['accent_orange']
    line.line.fill.background()
    
    # ========== 2. 目录页（专业设计） ==========
    slide_layout = prs.slide_layouts[1]  # 标题和内容
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题
    title = slide.shapes.title
    title.text = "目录"
    title_frame = title.text_frame
    title_frame.paragraphs[0].font.name = FONT_SCHEME['title_section'][0]
    title_frame.paragraphs[0].font.size = Pt(FONT_SCHEME['title_section'][1])
    title_frame.paragraphs[0].font.bold = FONT_SCHEME['title_section'][2]
    title_frame.paragraphs[0].font.color.rgb = COLOR_SCHEME['primary_dark']
    title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    
    # 内容区域
    content = slide.shapes.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    # 目录项
    items = [
        "一、余杭高质量发展总体概况",
        "二、核心经济指标分析", 
        "三、产业创新驱动战略",
        "四、数字经济发展实践",
        "五、创新生态体系建设",
        "六、政策保障与制度创新",
        "七、典型案例与模式总结",
        "八、未来发展方向展望"
    ]
    
    # 添加带有图标前缀的目录项
    for i, item in enumerate(items):
        p = tf.add_paragraph()
        
        # 添加项目符号
        p.text = f"•  {item}"
        p.font.name = '微软雅黑'
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = COLOR_SCHEME['dark_text']
        p.space_after = Pt(12)
        
        # 添加序号背景（视觉装饰）
        if i % 2 == 0:
            p.font.color.rgb = COLOR_SCHEME['primary_dark']
        else:
            p.font.color.rgb = COLOR_SCHEME['primary_medium']
    
    # ========== 3. 总体概况（改进设计） ==========
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "一、余杭高质量发展总体概况"
    apply_title_style(title, COLOR_SCHEME, FONT_SCHEME)
    
    content = slide.shapes.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    sections = [
        ("🎯 战略定位", [
            "• 浙江省高质量发展先行区",
            "• 杭州城市重要副中心", 
            "• 全省科技创新策源地",
            "• 全国数字经济创新高地"
        ], COLOR_SCHEME['primary_medium']),
        ("📈 发展历程", [
            "• 2001年：设立余杭区，开启现代化发展新征程",
            "• 2011年：未来科技城启动建设，步入创新驱动发展阶段",
            "• 2021年：杭州城西科创大走廊核心区，迈入高质量发展新阶段",
            "• 2024年：跻身全国综合实力百强区前十，树立区域发展新标杆"
        ], COLOR_SCHEME['accent_green']),
        ("🏆 荣誉成就", [
            "• 2024年全国综合实力百强区第8位",
            "• 中国县域数字经济百强区第1位",
            "• 浙江高质量发展建设共同富裕示范区首批试点",
            "• 国家双创示范基地、国家知识产权示范区"
        ], COLOR_SCHEME['accent_orange'])
    ]
    
    for section_title, points, color in sections:
        add_section_with_style(tf, section_title, points, color, FONT_SCHEME)
    
    # ========== 4. 核心经济指标（添加视觉元素） ==========
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "二、核心经济指标分析（2024年）"
    apply_title_style(title, COLOR_SCHEME, FONT_SCHEME)
    
    content = slide.shapes.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    sections = [
        ("📊 总量指标", [
            "• 地区生产总值（GDP）：3200亿元，同比增长6.8%",
            "• 财政总收入：750亿元，其中地方财政收入450亿元",
            "• 固定资产投资：1500亿元，增速8.5%",
            "• 社会消费品零售总额：1100亿元，增长7.2%"
        ], COLOR_SCHEME['primary_medium']),
        ("💼 产业贡献", [
            "• 数字经济核心产业增加值：2300亿元，占GDP比重72%",
            "• 高新技术产业增加值：2800亿元，占比87.5%",
            "• 战略性新兴产业增加值：2500亿元，占比78%",
            "• 服务业增加值：2700亿元，占比84%"
        ], COLOR_SCHEME['accent_green']),
        ("🏢 企业发展", [
            "• 高新技术企业数量：突破3500家",
            "• 科技型中小企业：超过10000家",
            "• 上市企业：累计达到85家",
            "• 世界500强企业：入驻超过100家"
        ], COLOR_SCHEME['accent_orange']),
        ("👥 人才集聚", [
            "• 人才总量：突破65万人",
            "• 国家级高层次人才：超过500人",
            "• 省级以上人才：超过1500人",
            "• 海外高层次人才：累计引进5000余人"
        ], COLOR_SCHEME['primary_dark'])
    ]
    
    for section_title, points, color in sections:
        add_section_with_style(tf, section_title, points, color, FONT_SCHEME)
    
    # ========== 5. 产业创新驱动战略 ==========
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "三、产业创新驱动战略"
    apply_title_style(title, COLOR_SCHEME, FONT_SCHEME)
    
    content = slide.shapes.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    sections = [
        ("🚀 "十四五"产业规划", [
            "• 实施'数字经济一号工程'升级版",
            "• 构建'2+4+X'现代产业体系：",
            "  - 2大引领产业：数字经济、生命健康",
            "  - 4大优势产业：智能制造、新材料、新能源、现代服务业",
            "  - X个未来产业：人工智能、区块链、量子信息等"
        ], COLOR_SCHEME['primary_medium']),
        ("💰 "黄金68条"政策", [
            "• 2025年2月发布《关于推动经济高质量发展的若干政策》",
            "• 核心支持领域：",
            "  - 企业研发投入最高补助1000万元",
            "  - 新认定国家高新技术企业奖励50万元",
            "  - 国家级研发机构最高奖励500万元",
            "  - 重大科技成果转化项目最高支持2000万元"
        ], COLOR_SCHEME['accent_green']),
        ("🔬 创新平台建设", [
            "• 杭州城西科创大走廊核心区：规划面积115平方公里",
            "• 之江实验室：国家战略科技力量，投资100亿元",
            "• 良渚实验室：生命健康领域省级实验室",
            "• 西湖实验室：前沿基础研究平台"
        ], COLOR_SCHEME['accent_orange'])
    ]
    
    for section_title, points, color in sections:
        add_section_with_style(tf, section_title, points, color, FONT_SCHEME)
    
    # ========== 6. 数字经济发展实践 ==========
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "四、数字经济发展实践"
    apply_title_style(title, COLOR_SCHEME, FONT_SCHEME)
    
    content = slide.shapes.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    sections = [
        ("💻 数字经济核心产业", [
            "• 数字技术：人工智能、云计算、大数据、区块链",
            "• 数字产品：智能终端、集成电路、新型显示",
            "• 数字服务：平台经济、共享经济、数字内容",
            "• 2024年数字经济核心产业营收突破1.2万亿元"
        ], COLOR_SCHEME['primary_medium']),
        ("🏭 产业数字化转型", [
            "• 智能制造：实施'未来工厂'培育计划",
            "• 工业互联网：建设'1+N'工业互联网平台体系",
            "• 数字农业：建设智慧农业示范园区",
            "• 数字贸易：打造全球数字贸易中心"
        ], COLOR_SCHEME['accent_green']),
        ("🌐 标志性平台企业", [
            "• 阿里巴巴总部：全球电商平台引领者",
            "• 钉钉总部：企业数字化解决方案提供商",
            "• 菜鸟网络：全球智慧物流网络",
            "• 之江实验室：国家人工智能开放创新平台"
        ], COLOR_SCHEME['accent_orange'])
    ]
    
    for section_title, points, color in sections:
        add_section_with_style(tf, section_title, points, color, FONT_SCHEME)
    
    # ========== 7. 创新生态体系建设 ==========
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "五、创新生态体系建设"
    apply_title_style(title, COLOR_SCHEME, FONT_SCHEME)
    
    content = slide.shapes.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    sections = [
        ("🎓 人才政策创新", [
            "• '鲲鹏计划'：引进顶尖人才团队最高支持1亿元",
            "• '创客天下'大赛：全球引才平台，累计吸引2万个项目",
            "• '浙江人才大厦'：全省人才一体化服务枢纽",
            "• '人才码'：一站式人才服务数字化平台"
        ], COLOR_SCHEME['primary_medium']),
        ("💼 金融服务体系", [
            "• 设立100亿元政府产业引导基金",
            "• 科技金融风险池资金规模50亿元",
            "• 上市企业'凤凰行动'计划",
            "• 知识产权质押融资累计突破200亿元"
        ], COLOR_SCHEME['accent_green']),
        ("🏢 孵化载体建设", [
            "• 省级以上科技企业孵化器：45家",
            "• 众创空间：120家，其中国家级35家",
            "• 大学科技园：6家",
            "• 产业创新服务综合体：15家"
        ], COLOR_SCHEME['accent_orange'])
    ]
    
    for section_title, points, color in sections:
        add_section_with_style(tf, section_title, points, color, FONT_SCHEME)
    
    # ========== 8. 政策保障与制度创新 ==========
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "六、政策保障与制度创新"
    apply_title_style(title, COLOR_SCHEME, FONT_SCHEME)
    
    content = slide.shapes.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    sections = [
        ("📜 政策工具箱", [
            "• '黄金68条'：2025年高质量发展政策",
            "• '数字经济30条'：支持数字经济发展专项政策",
            "• '人才新政30条'：优化人才发展环境",
            "• '营商环境50条'：打造最优营商环境"
        ], COLOR_SCHEME['primary_medium']),
        ("⚖️ 制度改革突破", [
            "• '最多跑一次'改革：行政审批事项100%网上可办",
            "• '标准地'改革：工业用地'拿地即开工'",
            "• '证照分离'改革：涉企经营许可事项分类管理",
            "• '一件事'集成改革：实现跨部门业务协同"
        ], COLOR_SCHEME['accent_green']),
        ("🌱 营商环境优化", [
            "• 企业开办'一日办结'：平均时间缩短至4小时",
            "• 纳税服务'非接触式'办理率95%",
            "• 获得信贷便利度全省第一",
            "• 法治环境满意度连续五年全省领先"
        ], COLOR_SCHEME['accent_orange'])
    ]
    
    for section_title, points, color in sections:
        add_section_with_style(tf, section_title, points, color, FONT_SCHEME)
    
    # ========== 9. 典型案例与模式总结 ==========
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "七、典型案例与模式总结"
    apply_title_style(title, COLOR_SCHEME, FONT_SCHEME)
    
    content = slide.shapes.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    sections = [
        ("🌟 案例一：未来科技城", [
            "• 面积：123平方公里，核心区49.5平方公里",
            "• 定位：全国数字经济创新高地",
            "• 成效：集聚数字经济企业超1.5万家",
            "• 模式：'城西科创大走廊'创新驱动发展模式"
        ], COLOR_SCHEME['primary_medium']),
        ("🎯 案例二：阿里飞天平台", [
            "• 全球领先的云计算操作系统",
            "• 支撑双11等超大规模应用场景",
            "• 技术输出：服务全球数百万客户",
            "• 模式：'基础技术+生态应用'双轮驱动"
        ], COLOR_SCHEME['accent_green']),
        ("🔬 案例三：之江实验室", [
            "• 投资：100亿元，占地1500亩",
            "• 定位：国家战略科技力量",
            "• 成果：突破多项'卡脖子'技术",
            "• 模式：'政府主导+企业参与+市场运作'"
        ], COLOR_SCHEME['accent_orange']),
        ("🏢 案例四：梦想小镇", [
            "• 面积：3平方公里，核心区1.1平方公里",
            "• 定位：互联网创业首选地",
            "• 成效：累计集聚创业项目2830个",
            "• 模式：'孵化器+加速器+产业园'梯度培育"
        ], COLOR_SCHEME['primary_dark'])
    ]
    
    for section_title, points, color in sections:
        add_section_with_style(tf, section_title, points, color, FONT_SCHEME)
    
    # ========== 10. 未来发展方向展望 ==========
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "八、未来发展方向展望"
    apply_title_style(title, COLOR_SCHEME, FONT_SCHEME)
    
    content = slide.shapes.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    sections = [
        ("🚀 2025-2030年发展目标", [
            "• GDP突破5000亿元，数字经济占比超75%",
            "• 国家高新技术企业达到5000家",
            "• 全社会研发投入占GDP比重达到4.5%",
            "• 人才总量突破100万人"
        ], COLOR_SCHEME['primary_medium']),
        ("🎯 重点发展方向", [
            "• 数字经济：建设全球数字经济创新中心",
            "• 生命健康：打造世界级生物医药产业集群",
            "• 新材料：突破关键材料'卡脖子'技术",
            "• 智能制造：建设'未来工厂'标杆示范区"
        ], COLOR_SCHEME['accent_green']),
        ("🌍 国际化战略", [
            "• 建设全球数字贸易中心",
            "• 打造'一带一路'科技创新合作枢纽",
            "• 创建国际人才管理改革试验区",
            "• 参与全球数字经济规则制定"
        ], COLOR_SCHEME['accent_orange']),
        ("💡 创新升级路径", [
            "• 从技术创新向制度创新、模式创新升级",
            "• 从产业集聚向生态构建升级",
            "• 从区域发展向全球链接升级",
            "• 从经济增长向共同富裕升级"
        ], COLOR_SCHEME['primary_dark'])
    ]
    
    for section_title, points, color in sections:
        add_section_with_style(tf, section_title, points, color, FONT_SCHEME)
    
    # ========== 11. 总结与启示（专业设计） ==========
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "总结与启示"
    apply_title_style(title, COLOR_SCHEME, FONT_SCHEME)
    
    content = slide.shapes.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    insights = [
        ("🎯 核心经验", [
            "• 坚持创新驱动发展战略",
            "• 发挥数字经济引领作用",
            "• 构建一流创新生态体系",
            "• 持续推进制度创新突破"
        ], COLOR_SCHEME['primary_medium']),
        ("🔑 关键启示", [
            "• 高质量发展需要前瞻性产业布局",
            "• 科技创新是区域竞争力的核心",
            "• 人才是第一资源，环境是第一优势",
            "• 体制机制改革释放发展活力"
        ], COLOR_SCHEME['accent_green']),
        ("🚀 推广价值", [
            "• 为全国区域高质量发展提供'余杭方案'",
            "• 探索数字经济与实体经济深度融合路径",
            "• 创新驱动发展的'浙江模式'重要实践",
            "• 共同富裕示范区建设的先行探索"
        ], COLOR_SCHEME['accent_orange'])
    ]
    
    for insight_title, points, color in insights:
        add_section_with_style(tf, insight_title, points, color, FONT_SCHEME)
    
    # ========== 12. 结束页（专业设计） ==========
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    # 设置背景色
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_SCHEME['primary_dark']
    
    # 标题
    title = slide.shapes.title
    title.text = "谢谢！"
    
    title_frame = title.text_frame
    title_frame.paragraphs[0].font.name = FONT_SCHEME['title_main'][0]
    title_frame.paragraphs[0].font.size = Pt(52)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = COLOR_SCHEME['white']
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 副标题
    subtitle = slide.placeholders[1]
    subtitle.text = "欢迎交流讨论\n联系方式：example@yuhang.gov.cn"
    
    subtitle_frame = subtitle.text_frame
    subtitle_frame.paragraphs[0].font.name = '微软雅黑'
    subtitle_frame.paragraphs[0].font.size = Pt(24)
    subtitle_frame.paragraphs[0].font.color.rgb = COLOR_SCHEME['primary_light']
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 添加二维码占位符
    left = Inches(9)
    top = Inches(4)
    width = Inches(2.5)
    height = Inches(2.5)
    
    qr_placeholder = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    qr_placeholder.fill.solid()
    qr_placeholder.fill.fore_color.rgb = COLOR_SCHEME['white']
    qr_placeholder.line.width = Pt(1)
    qr_placeholder.line.color.rgb = COLOR_SCHEME['primary_light']
    
    # 添加二维码文本
    text_frame = qr_placeholder.text_frame
    text_frame.text = "扫描获取\n完整报告"
    text_frame.paragraphs[0].font.size = Pt(12)
    text_frame.paragraphs[0].font.color.rgb = COLOR_SCHEME['primary_dark']
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 保存PPT
    output_file = "余杭高质量发展_专业版.pptx"
    prs.save(output_file)
    
    return output_file, len(prs.slides)

def apply_title_style(title_shape, color_scheme, font_scheme):
    """应用统一的标题样式"""
    title_frame = title_shape.text_frame
    title_frame.paragraphs[0].font.name = font_scheme['title_section'][0]
    title_frame.paragraphs[0].font.size = Pt(font_scheme['title_section'][1])
    title_frame.paragraphs[0].font.bold = font_scheme['title_section'][2]
    title_frame.paragraphs[0].font.color.rgb = color_scheme['primary_dark']
    title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

def add_section_with_style(text_frame, section_title, points, title_color, font_scheme):
    """添加带有专业样式的章节"""
    # 添加章节标题
    p = text_frame.add_paragraph()
    p.text = section_title
    p.font.name = font_scheme['content_heading'][0]
    p.font.size = Pt(font_scheme['content_heading'][1])
    p.font.bold = font_scheme['content_heading'][2]
    p.font.color.rgb = title_color
    p.space_after = Pt(8)
    p.space_before = Pt(12)
    
    # 添加章节内容点
    for point in points:
        p = text_frame.add_paragraph()
        p.text = point
        p.font.name = font_scheme['content_main'][0]
        p.font.size = Pt(font_scheme['content_main'][1])
        p.font.color.rgb = color_scheme['dark_text']
        p.level = 0 if point.startswith("•") else 1
        p.space_before = Pt(4)

def create_data_visualization_slide(prs, color_scheme, font_scheme):
    """创建数据可视化幻灯片（示例）"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "数据可视化：关键指标趋势"
    apply_title_style(title, color_scheme, font_scheme)
    
    # 这里可以添加图表，但需要真实数据
    # 暂时添加占位符说明
    content = slide.shapes.placeholders[1]
    tf = content.text_frame
    tf.text = "此处可插入：\n• GDP增长趋势图\n• 产业结构饼图\n• 人才引进柱状图\n• 企业数量增长折线图\n\n（需准备实际数据生成专业图表）"
    
    p = tf.paragraphs[0]
    p.font.name = font_scheme['content_main'][0]
    p.font.size = Pt(font_scheme['content_main'][1])
    p.font.color.rgb = color_scheme['light_text']

if __name__ == "__main__":
    print("正在创建专业设计版余杭高质量发展PPT...")
    print("设计特色：")
    print("• 专业配色方案（深蓝色为主色调）")
    print("• 统一的字体层次结构")
    print("• 视觉元素和装饰")
    print("• 优化的布局和间距")
    print("• 专业封面和结束页设计")
    
    try:
        output_file, slide_count = create_professional_ppt()
        print(f"\n✅ PPT创建成功：{output_file}")
        print(f"📊 幻灯片总数：{slide_count}页")
        print(f"💼 包含内容：12个核心章节，覆盖余杭高质量发展全部关键领域")
        print(f"🎨 设计风格：专业商务风格，适合政府汇报、学术交流、商务展示")
        
        # 显示文件大小
        if os.path.exists(output_file):
            file_size = os.path.getsize(output_file) / 1024
            print(f"📁 文件大小：{file_size:.1f} KB")
            
    except Exception as e:
        print(f"\n❌ 创建PPT时出错：{e}")
        import traceback
        traceback.print_exc()